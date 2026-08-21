"""Zet de presentatornotitie van een slide: de toelichting die de spreker zegt.

Usage:
    python set_notes.py <unpacked_dir> --json '{"3": "wat de spreker zegt", "4": [...]}'
    python set_notes.py <deck.pptx>    --json '{"3": "..."}'
    python set_notes.py <unpacked_dir> --file notities.json
    python set_notes.py <unpacked_dir> --slide 3 --text "wat de spreker zegt"
    python set_notes.py <deck.pptx> --check          # alleen rapporteren

Waarom dit script bestaat
-------------------------
Het dichtheidsbesluit (SKILL.md, stap 1, besluit 1) zegt bij een spreekdeck dat wat de
spreker zegt níét op de slide hoort en dat de toelichting naar de presentatornotities
gaat. Tot nu toe schreef geen enkel script een notesSlide, dus was dat besluit half
uitvoerbaar: je kon de toelichting van de slide halen maar niet in het bestand bewaren,
en ze bleef in de outline achter. Wie het deck opent ziet dan een kale slide zonder de
tekst die erbij hoort.

De sleutel is de DECKPOSITIE — het nummer dat de lezer ziet, hetzelfde nummer dat
`fit_title.py` en `qa_tellingen.py` in hun bevindingen zetten, niet het bestandsnummer.
Die twee lopen uiteen zodra er met `add_slide.py --at` iets tussen is gezet.

De waarde is een string (één alinea) of een lijst strings (meer alinea's):

    {"3": "Hier vertel je dat de doorlooptijd met 39 dagen daalde."}
    {"3": ["Eerste alinea.", "Tweede alinea."]}
    {"3": ""}      leegmaken

Uitgepakt of ingepakt
---------------------
Beide, en het script kijkt zelf wat je hem geeft: een map met `ppt/slides` erin is de
uitgepakte boom, een `.pptx` is het ingepakte bestand. Dat verschil is hier wél van
belang, anders dan bij `add_chart.py` en `add_table.py` die alleen ná `pack.py` werken —
en waarvan een tweede `pack.py` het werk dus sloopt. Notities horen bij de tekst en niet
bij de grafiek, dus de normale plek is de UITGEPAKTE boom: dan overleven ze `clean.py`
en `pack.py` gewoon mee. De ingepakte route is er voor een deck dat al opgeleverd is en
waar alleen de notities nog bij moeten.

Herhaalbaar zonder te verdubbelen: een tweede aanroep VERVANGT de notitie van die slide,
precies zoals `set_text.py` de alinea's van een placeholder vervangt. Er komt nooit een
tweede notesSlide bij dezelfde slide, want een bestaande wordt hergebruikt.

Een lege waarde maakt de notitie leeg; `clean.py` gooit de lege notesSlide daarna weg.
Dat is dezelfde regel als voor placeholders: wat je leeg laat verdwijnt.

Font, maat en kleur komen uit de notesMaster — dat is layout-first, net als op de slide.
Rechte apostrofs worden typografisch (`risico's` → `risico’s`).

Output is compacte JSON.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from lxml import etree

from _deck import deck_positions, emit, normalise_apostrophes

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
NS = {"a": A, "p": P, "r": R}

NOTES_SLIDE_REL = f"{R}/notesSlide"
NOTES_MASTER_REL = f"{R}/notesMaster"
SLIDE_REL = f"{R}/slide"
NOTES_CT = (
    "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"
)

#: De placeholders die een notesSlide van de notesMaster overneemt: de dia-afbeelding,
#: het notitievak en het paginanummer. Kop-, datum- en voettekstplaceholders niet — die
#: laat PowerPoint zelf ook weg bij een nieuwe notitiepagina, en python-pptx doet in
#: `clone_master_placeholders()` precies deze drie.
CLONED_PH = ("sldImg", "body", "sldNum")


def paragrafen(value) -> list[str]:
    """De waarde uit de JSON naar een lijst alinea's."""
    if value is None:
        return []
    if isinstance(value, str):
        tekst = value.strip()
        return [tekst] if tekst else []
    if isinstance(value, list):
        regels = []
        for item in value:
            if not isinstance(item, str):
                raise SystemExit(
                    f"een notitie is tekst, geen {type(item).__name__}: {item!r} — "
                    "geef een string of een lijst strings"
                )
            if item.strip():
                regels.append(item.strip())
        return regels
    raise SystemExit(
        f"kan notitie niet lezen: {value!r} — geef een string of een lijst strings"
    )


def net(regels: list[str]) -> list[str]:
    return [normalise_apostrophes(regel) for regel in regels]


# ------------------------------------------------------------------ uitgepakt


def _rels_path(part: Path) -> Path:
    return part.parent / "_rels" / f"{part.name}.rels"


def _lees_rels(path: Path) -> etree._ElementTree | None:
    return etree.parse(str(path)) if path.exists() else None


def _schrijf(tree: etree._ElementTree, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    tree.write(str(path), xml_declaration=True, encoding="UTF-8", standalone=True)


def _nieuwe_rid(root: etree._Element) -> str:
    gebruikt = {
        int(m.group(1))
        for rel in root
        for m in [re.fullmatch(r"rId(\d+)", rel.get("Id", ""))]
        if m
    }
    return f"rId{max(gebruikt, default=0) + 1}"


def notes_part_van(slide_path: Path, unpacked: Path) -> Path | None:
    """Het notesSlide-bestand waar deze slide al naar verwijst, of None."""
    tree = _lees_rels(_rels_path(slide_path))
    if tree is None:
        return None
    for rel in tree.getroot():
        if rel.get("Type") == NOTES_SLIDE_REL:
            doel = rel.get("Target", "")
            return (slide_path.parent / doel).resolve()
    return None


def _leeg_sp(id_: int, naam: str, ph_type: str, idx: str | None, sz: str | None,
            met_txbody: bool) -> etree._Element:
    """Eén lege placeholder-shape, zoals python-pptx er een op een notesSlide zet."""
    sp = etree.Element(f"{{{P}}}sp")
    nv = etree.SubElement(sp, f"{{{P}}}nvSpPr")
    cnv = etree.SubElement(nv, f"{{{P}}}cNvPr")
    cnv.set("id", str(id_))
    cnv.set("name", naam)
    cnv_sp = etree.SubElement(nv, f"{{{P}}}cNvSpPr")
    etree.SubElement(cnv_sp, f"{{{A}}}spLocks").set("noGrp", "1")
    nv_pr = etree.SubElement(nv, f"{{{P}}}nvPr")
    ph = etree.SubElement(nv_pr, f"{{{P}}}ph")
    ph.set("type", ph_type)
    if sz:
        ph.set("sz", sz)
    if idx:
        ph.set("idx", idx)
    etree.SubElement(sp, f"{{{P}}}spPr")
    if met_txbody:
        body = etree.SubElement(sp, f"{{{P}}}txBody")
        etree.SubElement(body, f"{{{A}}}bodyPr")
        etree.SubElement(body, f"{{{A}}}lstStyle")
        etree.SubElement(body, f"{{{A}}}p")
    return sp


def _nieuwe_notesslide_xml(unpacked: Path) -> etree._ElementTree:
    """Een lege notesSlide met de placeholders van de notesMaster erop."""
    master = unpacked / "ppt" / "notesMasters" / "notesMaster1.xml"
    root = etree.Element(f"{{{P}}}notes", nsmap={"a": A, "p": P, "r": R})
    csld = etree.SubElement(root, f"{{{P}}}cSld")
    tree_el = etree.SubElement(csld, f"{{{P}}}spTree")
    nv_grp = etree.SubElement(tree_el, f"{{{P}}}nvGrpSpPr")
    cnv = etree.SubElement(nv_grp, f"{{{P}}}cNvPr")
    cnv.set("id", "1")
    cnv.set("name", "")
    etree.SubElement(nv_grp, f"{{{P}}}cNvGrpSpPr")
    etree.SubElement(nv_grp, f"{{{P}}}nvPr")
    grp_pr = etree.SubElement(tree_el, f"{{{P}}}grpSpPr")
    xfrm = etree.SubElement(grp_pr, f"{{{A}}}xfrm")
    for tag, a, b in (("off", "x", "y"), ("ext", "cx", "cy"),
                      ("chOff", "x", "y"), ("chExt", "cx", "cy")):
        el = etree.SubElement(xfrm, f"{{{A}}}{tag}")
        el.set(a, "0")
        el.set(b, "0")

    volgend = 2
    if master.exists():
        master_tree = etree.parse(str(master))
        for sp in master_tree.getroot().iter(f"{{{P}}}sp"):
            ph = sp.find(f".//{{{P}}}nvPr/{{{P}}}ph")
            if ph is None or ph.get("type") not in CLONED_PH:
                continue
            ph_type = ph.get("type")
            naam = {"sldImg": "Slide Image Placeholder",
                    "body": "Notes Placeholder",
                    "sldNum": "Slide Number Placeholder"}[ph_type]
            tree_el.append(
                _leeg_sp(volgend, f"{naam} {volgend}", ph_type, ph.get("idx"),
                         ph.get("sz"), met_txbody=ph_type != "sldImg")
            )
            volgend += 1
    if not any(
        sp.find(f".//{{{P}}}nvPr/{{{P}}}ph").get("type") == "body"
        for sp in tree_el.iter(f"{{{P}}}sp")
    ):
        # Geen notesMaster in de boom: dan nog steeds een notitievak, anders is er
        # niets om de tekst in te zetten.
        tree_el.append(_leeg_sp(volgend, "Notes Placeholder", "body", "3",
                                "quarter", met_txbody=True))

    ovr = etree.SubElement(root, f"{{{P}}}clrMapOvr")
    etree.SubElement(ovr, f"{{{A}}}masterClrMapping")
    return etree.ElementTree(root)


def _volgend_notesnummer(unpacked: Path) -> int:
    notes_dir = unpacked / "ppt" / "notesSlides"
    gebruikt = {
        int(m.group(1))
        for path in notes_dir.glob("notesSlide*.xml")
        for m in [re.fullmatch(r"notesSlide(\d+)\.xml", path.name)]
        if m
    }
    return max(gebruikt, default=0) + 1


def _voeg_override_toe(unpacked: Path, partname: str) -> None:
    ct_path = unpacked / "[Content_Types].xml"
    tree = etree.parse(str(ct_path))
    root = tree.getroot()
    ns = root.nsmap.get(None)
    tag = f"{{{ns}}}Override" if ns else "Override"
    for override in root.findall(tag):
        if override.get("PartName") == partname:
            return
    override = etree.SubElement(root, tag)
    override.set("PartName", partname)
    override.set("ContentType", NOTES_CT)
    _schrijf(tree, ct_path)


def _maak_notes_part(unpacked: Path, slide_path: Path) -> Path:
    """Nieuw notesSlide-bestand, met zijn rels, zijn override en de rel op de slide."""
    if not (unpacked / "ppt" / "notesMasters" / "notesMaster1.xml").exists():
        print(
            "let op: er staat geen notesMaster in deze boom; de notitie krijgt dan geen "
            "opmaak uit het sjabloon",
            file=sys.stderr,
        )
    notes_dir = unpacked / "ppt" / "notesSlides"
    notes_dir.mkdir(parents=True, exist_ok=True)
    nummer = _volgend_notesnummer(unpacked)
    notes_path = notes_dir / f"notesSlide{nummer}.xml"
    _schrijf(_nieuwe_notesslide_xml(unpacked), notes_path)

    rels_root = etree.Element(f"{{{PKG_REL}}}Relationships", nsmap={None: PKG_REL})
    rel1 = etree.SubElement(rels_root, f"{{{PKG_REL}}}Relationship")
    rel1.set("Id", "rId1")
    rel1.set("Type", NOTES_MASTER_REL)
    rel1.set("Target", "../notesMasters/notesMaster1.xml")
    rel2 = etree.SubElement(rels_root, f"{{{PKG_REL}}}Relationship")
    rel2.set("Id", "rId2")
    rel2.set("Type", SLIDE_REL)
    rel2.set("Target", f"../slides/{slide_path.name}")
    _schrijf(etree.ElementTree(rels_root), _rels_path(notes_path))

    slide_rels_path = _rels_path(slide_path)
    slide_rels = _lees_rels(slide_rels_path)
    if slide_rels is None:
        root = etree.Element(f"{{{PKG_REL}}}Relationships", nsmap={None: PKG_REL})
        slide_rels = etree.ElementTree(root)
    root = slide_rels.getroot()
    rel = etree.SubElement(root, f"{{{PKG_REL}}}Relationship")
    rel.set("Id", _nieuwe_rid(root))
    rel.set("Type", NOTES_SLIDE_REL)
    rel.set("Target", f"../notesSlides/{notes_path.name}")
    _schrijf(slide_rels, slide_rels_path)

    _voeg_override_toe(unpacked, f"/ppt/notesSlides/{notes_path.name}")
    return notes_path


def _notitievak(notes_root: etree._Element) -> etree._Element | None:
    for sp in notes_root.iter(f"{{{P}}}sp"):
        ph = sp.find(f".//{{{P}}}nvPr/{{{P}}}ph")
        if ph is not None and ph.get("type") == "body":
            return sp
    return None


def _schrijf_alineas(sp: etree._Element, regels: list[str]) -> None:
    body = sp.find(f"{{{P}}}txBody")
    if body is None:
        raise SystemExit("het notitievak heeft geen txBody")
    for para in body.findall(f"{{{A}}}p"):
        body.remove(para)
    for regel in regels or [""]:
        node = etree.SubElement(body, f"{{{A}}}p")
        if not regel:
            continue
        r = etree.SubElement(node, f"{{{A}}}r")
        rpr = etree.SubElement(r, f"{{{A}}}rPr")
        rpr.set("lang", "nl-NL")
        rpr.set("dirty", "0")
        etree.SubElement(r, f"{{{A}}}t").text = regel


def tekst_van(notes_path: Path) -> list[str]:
    if not notes_path.exists():
        return []
    sp = _notitievak(etree.parse(str(notes_path)).getroot())
    if sp is None:
        return []
    regels = []
    for para in sp.iter(f"{{{A}}}p"):
        regel = "".join(t.text or "" for t in para.iter(f"{{{A}}}t")).strip()
        if regel:
            regels.append(regel)
    return regels


def op_uitgepakt(unpacked: Path, gewenst: dict[int, list[str]], check: bool) -> dict:
    posities = deck_positions(unpacked)
    slides_dir = unpacked / "ppt" / "slides"
    per_positie = {
        positie: slides_dir / naam for naam, positie in posities.items()
    }
    onbekend = sorted(set(gewenst) - set(per_positie))
    if onbekend:
        raise SystemExit(
            f"deze deckposities bestaan niet: {onbekend} — het deck heeft "
            f"{len(per_positie)} slides"
        )

    gezet, geleegd, aanwezig = [], [], []
    for positie in sorted(per_positie):
        slide_path = per_positie[positie]
        notes_path = notes_part_van(slide_path, unpacked)
        if positie in gewenst and not check:
            regels = net(gewenst[positie])
            if notes_path is None:
                if not regels:
                    continue
                notes_path = _maak_notes_part(unpacked, slide_path)
            tree = etree.parse(str(notes_path))
            sp = _notitievak(tree.getroot())
            if sp is None:
                raise SystemExit(
                    f"{notes_path.name} heeft geen notitievak — verwijder het bestand "
                    "en de rel op de slide en laat dit script hem opnieuw aanmaken"
                )
            _schrijf_alineas(sp, regels)
            _schrijf(tree, notes_path)
            (gezet if regels else geleegd).append(positie)
        if notes_path is not None and notes_path.exists():
            regels = tekst_van(notes_path)
            aanwezig.append(
                {
                    "slide": positie,
                    "file": slide_path.name,
                    "notes": notes_path.name,
                    "paragrafen": len(regels),
                    "tekens": sum(len(r) for r in regels),
                }
            )
    return {
        "doel": str(unpacked),
        "vorm": "uitgepakt",
        "numbering": "slide = deckpositie, file = bestandsnaam",
        "gezet": gezet,
        "geleegd": geleegd,
        "notities": aanwezig,
    }


# ------------------------------------------------------------------ ingepakt


def op_pptx(deck: Path, gewenst: dict[int, list[str]], check: bool) -> dict:
    from pptx import Presentation

    prs = Presentation(str(deck))
    slides = list(prs.slides)
    onbekend = sorted(p for p in gewenst if not 1 <= p <= len(slides))
    if onbekend:
        raise SystemExit(
            f"deze deckposities bestaan niet: {onbekend} — het deck heeft "
            f"{len(slides)} slides"
        )

    gezet, geleegd, aanwezig = [], [], []
    for positie, slide in enumerate(slides, start=1):
        if positie in gewenst and not check:
            regels = net(gewenst[positie])
            frame = slide.notes_slide.notes_text_frame
            frame.clear()
            if regels:
                frame.paragraphs[0].text = regels[0]
                for regel in regels[1:]:
                    frame.add_paragraph().text = regel
            (gezet if regels else geleegd).append(positie)
        if not slide.has_notes_slide:
            continue
        frame = slide.notes_slide.notes_text_frame
        regels = [p.text.strip() for p in frame.paragraphs if p.text.strip()]
        aanwezig.append(
            {
                "slide": positie,
                "notes": slide.notes_slide.part.partname.split("/")[-1],
                "paragrafen": len(regels),
                "tekens": sum(len(r) for r in regels),
            }
        )
    if not check and (gezet or geleegd):
        prs.save(str(deck))
    return {
        "doel": str(deck),
        "vorm": "ingepakt",
        "numbering": "slide = deckpositie",
        "gezet": gezet,
        "geleegd": geleegd,
        "notities": aanwezig,
    }


# ------------------------------------------------------------------ cli


def lees_payload(args) -> dict[int, list[str]]:
    if args.slide is not None:
        if args.text is None:
            raise SystemExit("--slide vraagt --text")
        return {args.slide: paragrafen(args.text)}
    if not args.payload and not args.file:
        if args.check:
            return {}
        raise SystemExit("geef --json, --file, of --slide met --text (of --check)")
    raw = args.file.read_text(encoding="utf-8") if args.file else args.payload
    inhoud = json.loads(raw)
    if not isinstance(inhoud, dict):
        raise SystemExit(
            'de JSON is een object met deckposities als key: {"3": "toelichting"}'
        )
    gewenst = {}
    for key, value in inhoud.items():
        try:
            positie = int(key)
        except (TypeError, ValueError):
            raise SystemExit(
                f"onbekende slidekey {key!r} — de key is de DECKPOSITIE, het nummer dat "
                "de lezer ziet"
            ) from None
        gewenst[positie] = paragrafen(value)
    return gewenst


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("doel", type=Path, help="uitgepakte map of .pptx")
    parser.add_argument("--json", dest="payload", help="JSON: deckpositie -> notitie")
    parser.add_argument("--file", type=Path, help="lees die JSON uit een bestand")
    parser.add_argument("--slide", type=int, help="één deckpositie")
    parser.add_argument("--text", help="de notitie bij --slide")
    parser.add_argument("--check", action="store_true", help="alleen rapporteren")
    args = parser.parse_args()

    gewenst = lees_payload(args)

    if args.doel.is_dir():
        if not (args.doel / "ppt" / "slides").is_dir():
            raise SystemExit(f"geen ppt/slides in {args.doel}")
        emit(op_uitgepakt(args.doel, gewenst, args.check))
        return
    if args.doel.suffix.lower() not in {".pptx", ".potx"}:
        raise SystemExit(
            f"{args.doel} is geen uitgepakte map en geen .pptx — geef de builddir "
            "(dan overleven de notities clean.py en pack.py) of het ingepakte deck"
        )
    if not args.doel.exists():
        raise SystemExit(f"{args.doel} bestaat niet")
    emit(op_pptx(args.doel, gewenst, args.check))


if __name__ == "__main__":
    if len(sys.argv) == 1:
        print(__doc__)
        sys.exit(1)
    main()
