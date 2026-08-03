"""Hygiënecontrole op een gebouwde deck. Geen vormbeoordeling.

Usage:
    python qa_text.py <deck.pptx>
    python qa_text.py <deck.pptx> --strict     # ook waarschuwingen geven exit 1

Wat dit script wél doet: kijken of er iets in de deck staat wat er niet in mag omdat het
een restje, een fout of een leesbaarheidsprobleem is. Een placeholder die nooit gevuld is.
Een `{{MARKER}}` uit het concept. Calibri, omdat iemand een `<a:latin/>` vergat. Een
hardgecodeerde hex in plaats van `schemeClr`. Een rechte apostrof. Een titel in
onderkast.

Wat dit script NIET doet: beoordelen of de slide mooi is, of de vakken vol genoeg staan,
hoeveel tekstgroottes er op een slide staan, of de compositie klopt. Dat oordeel komt van
de render, niet van een meting. De vorige versie van dit script deed dat wel, samen met
`qa_fit.py` en `qa_typography.py`, en het resultaat was dat de bouwer regels ging vermijden
in plaats van slides ging maken.

Twee uitzonderingen op die grens, en die zijn er omdat ze in de praktijk misgingen. De
titelletter in de contentzone is een `critical`: Gotham Bold hoort in de titel en komt daar
uit de layout, en op de slide zelf is de letter Montserrat Light of Lato Light. En de
frequentie van de drager is een `warn`: staat er op elke slide letter van 28pt of groter,
dan trekt die maat geen aandacht meer. Beide zijn hard te tellen en geen smaakoordeel.

De bevindingen met severity `critical` betekenen allemaal "dit mag niet naar een klant"
en geen ervan gaat over vormgeving. Daar hoort ook autofit bij: een vak waarin PowerPoint
de tekst stil mag krimpen (`normAutofit`, of een placeholder die hem uit de layout erft)
is een defect, want het font gaat nooit omlaag om passend te worden. Liever tekst die
zichtbaar te lang is en een mens die beslist, dan een vak dat zichzelf stilletjes
verkleint. De rest is `warn`: het kan kloppen, kijk er even naar.

Output is compacte JSON. Exit 1 zodra er een `critical` staat, of bij `--strict` ook op een
`warn`.
"""

from __future__ import annotations

import argparse
import re
import sys
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402

from _deck import BRAND_FONTS, emit  # noqa: E402
from shapes import DRAGER_PLAFOND, DRAGER_VLOER, TITELFONT  # noqa: E402

# `{{IETS}}` uit een concept, en een niet-gesloten `{{` die de eerste niet vangt.
MARKER = re.compile(r"\{\{[A-Z0-9 _-]{2,40}\}\}")
RESIDUE = re.compile(r"\{\{|\}\}|<<[A-Z]|\bTODO\b|\bTBD\b|\bXXX\b|\[invullen\]", re.I)

# Sjabloonprompts die als SUBSTRING altijd fout zijn.
TEMPLATE_PROMPTS = (
    "klik om tekst toe te voegen",
    "klik hier om",
    "voer titel in",
    "titel toevoegen",
    "lorem ipsum",
    "tekst hier",
    "your text here",
    "master title style",
    "master text styles",
)

# Sjabloonwoorden die alleen fout zijn als ze de HELE tekst van een vak zijn: "titel" en
# "tekst" komen in gewone Nederlandse slidetekst voortdurend voor.
TEMPLATE_PROMPTS_EXACT = {
    "presentatie", "titel", "subtitel", "ondertitel", "kop", "koptekst",
    "tekst", "titeltekst", "naam", "datum",
}

# Chrome uit de layout: paginanummer, datum, voettekst. Geen inhoud, dus niet toetsen.
CHROME_TYPES = {"SLIDE_NUMBER", "DATE", "FOOTER", "HEADER"}

STRAIGHT_QUOTE = re.compile(r"[\w)](')[\w]|(\")")
DOUBLE_SPACE = re.compile(r"\S(  +)\S")


def add(findings, slide, check, message, severity="warn"):
    findings.append(
        {"slide": slide, "check": check, "severity": severity, "message": message}
    )


def is_caps(text: str) -> bool:
    letters = [c for c in text if c.isalpha()]
    return all(c.isupper() for c in letters) if letters else True


def is_chrome(shape) -> bool:
    try:
        return str(shape.placeholder_format.type).split()[0] in CHROME_TYPES
    except (AttributeError, ValueError):
        return False


# Layouts zonder `type="title"`, waar idx 10 of 14 de kop draagt: de covers en de
# sectiedividers. Op layout 20 is idx 10 juist de CONTENTZONE, dus de idx-regel mag daar
# niet gelden — anders wordt elke eerste alinea van een tekstslide als kop beoordeeld en
# afgekeurd omdat hij niet in kapitalen staat.
HEADLINE_BY_IDX_LAYOUTS = {1, 4} | set(range(6, 17))

# De oranje outro (layout 2 en 3) draagt logo en kleurvlak in de LAYOUT en heeft op de
# slide zelf nul vormen. Dat is geen lege slide maar een afsluiter zoals hij bedoeld is.
# Layout 17 staat hier NIET bij: een blanco canvas zonder vormen is wél een defect.
INTENTIONALLY_BLANK_LAYOUTS = {2, 3}


def layout_number(slide) -> int | None:
    match = re.search(r"slideLayout(\d+)\.xml", str(slide.slide_layout.part.partname))
    return int(match.group(1)) if match else None


def is_title(shape, layout: int | None) -> bool:
    """De kop van de slide: een echte titelplaceholder, of idx 10/14 op cover en divider."""
    try:
        fmt = shape.placeholder_format
    except (AttributeError, ValueError):
        return False
    if str(fmt.type).split()[0] in {"TITLE", "CENTER_TITLE"}:
        return True
    return layout in HEADLINE_BY_IDX_LAYOUTS and fmt.idx in {10, 14}


def walk(shapes):
    """Alle vormen, groepen doorlopen."""
    for shape in shapes:
        if getattr(shape, "shape_type", None) is not None and shape.shape_type == 6:
            yield from walk(shape.shapes)
            continue
        yield shape


def runs_of(slide):
    """(shape, run) voor elke tekstrun op de slide, groepen doorlopen."""
    for shape in walk(slide.shapes):
        if not getattr(shape, "has_text_frame", False) or is_chrome(shape):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                yield shape, run


def autofit_status(shape) -> str | None:
    """Hoe deze vorm met autofit omgaat: 'normAutofit'/'spAutoFit' als PowerPoint de
    tekst mag krimpen of het vak mag oprekken, 'erft' als een placeholder geen eigen
    keuze draagt en dus de layout volgt (en dit sjabloon zet daar normAutofit), en
    None als er een expliciete <a:noAutofit/> staat of een eigen vorm geen keuze
    nodig heeft (zonder element doet OOXML niets)."""
    from pptx.oxml.ns import qn
    body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        for tag in ("a:normAutofit", "a:spAutoFit"):
            if body_pr.find(qn(tag)) is not None:
                return tag.split(":")[1]
        if body_pr.find(qn("a:noAutofit")) is not None:
            return None
    return "erft" if getattr(shape, "is_placeholder", False) else None


def analyse(deck: Path) -> dict:
    presentation = Presentation(str(deck))
    findings: list[dict] = []
    fonts: Counter = Counter()
    charts = tables = 0
    contentslides = 0
    slides_met_drager: list[int] = []

    for number, slide in enumerate(presentation.slides, start=1):
        has_text = False
        reported: set[tuple] = set()
        layout = layout_number(slide)
        titels = {id(s) for s in walk(slide.shapes) if is_title(s, layout)}
        if layout not in HEADLINE_BY_IDX_LAYOUTS | INTENTIONALLY_BLANK_LAYOUTS:
            contentslides += 1

        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                charts += 1
            if getattr(shape, "has_table", False):
                tables += 1

        for shape in walk(slide.shapes):
            if not getattr(shape, "has_text_frame", False) or is_chrome(shape):
                continue
            status = autofit_status(shape)
            if status is None:
                continue
            naam = getattr(shape, "name", "?")
            if status == "erft":
                add(findings, number, "autofit",
                    f'"{naam}" draagt geen <a:noAutofit/> en erft dus de autofit van '
                    "de layout, en dit sjabloon zet daar normAutofit. Zet noAutofit "
                    "expliciet.", "critical")
            else:
                add(findings, number, "autofit",
                    f'"{naam}" staat op {status}: PowerPoint mag de tekst stil '
                    "krimpen of het vak oprekken. Zet <a:noAutofit/>; past de tekst "
                    "dan niet, dan is dat zichtbaar op de render en beslist een mens "
                    "— korter of een groter vak, nooit een kleiner font.", "critical")

        for shape, run in runs_of(slide):
            text = run.text
            if text.strip():
                has_text = True

            for hit in MARKER.findall(text):
                add(findings, number, "marker",
                    f"{hit} staat er nog in — vul de placeholder of haal hem weg",
                    "critical")
            if not MARKER.search(text):
                for match in RESIDUE.finditer(text):
                    add(findings, number, "placeholder-residue",
                        f'"{match.group(0)}" is een restje uit het concept: '
                        f'"{text.strip()[:60]}"',
                        "critical")
                    break

            lowered = text.lower()
            for prompt in TEMPLATE_PROMPTS:
                if prompt in lowered:
                    add(findings, number, "marker",
                        f'sjabloonprompt op de slide: "{text.strip()[:60]}"', "critical")
                    break

            in_titel = id(shape) in titels
            pt = run.font.size.pt if run.font.size is not None else None

            font = run.font.name
            if font:
                fonts[font] += 1
                if (font.strip().lower() == TITELFONT.lower() and not in_titel
                        and ("titelfont", 0) not in reported):
                    reported.add(("titelfont", 0))
                    add(findings, number, "titelfont-in-body",
                        f'"{getattr(shape, "name", "?")}" schrijft {TITELFONT} in de '
                        "contentzone. Die letter hoort alleen in de titel en komt daar uit "
                        "de layout; op de slide zelf is het Montserrat Light of Lato Light.",
                        "critical")
                if font not in BRAND_FONTS and ("font", font) not in reported:
                    reported.add(("font", font))
                    add(findings, number, "off-brand-font",
                        f'"{font}" is geen huisstijlfont. Meestal betekent dit dat een '
                        "eigen vorm geen expliciete <a:latin/> heeft en Calibri erft.")

            try:
                rgb = run.font.color.rgb
            except (AttributeError, TypeError):
                rgb = None
            if rgb is not None and ("color", str(rgb)) not in reported:
                reported.add(("color", str(rgb)))
                add(findings, number, "off-brand-color",
                    f"#{rgb} is een harde hex — gebruik schemeClr, anders drijft de "
                    "kleur weg van het thema.")

            if pt is not None and not in_titel and text.strip():
                if pt >= DRAGER_VLOER and number not in slides_met_drager:
                    slides_met_drager.append(number)
                if pt > DRAGER_PLAFOND and ("luid", 0) not in reported:
                    reported.add(("luid", 0))
                    add(findings, number, "drager-te-groot",
                        f"{pt:.0f}pt in de contentzone staat boven de band van "
                        f"{DRAGER_VLOER:.0f} tot {DRAGER_PLAFOND:.0f}pt: dan neemt de "
                        f'aandachtstrekker de slide over. "{text.strip()[:40]}"')

            if STRAIGHT_QUOTE.search(text):
                add(findings, number, "straight-quote",
                    f'rechte apostrof of aanhalingsteken: "{text.strip()[:60]}"')
            if DOUBLE_SPACE.search(text):
                add(findings, number, "double-space",
                    f'dubbele spatie: "{text.strip()[:60]}"')

        for shape in slide.shapes:
            if not is_title(shape, layout) or not getattr(shape, "has_text_frame", False):
                continue
            title = shape.text_frame.text.strip()
            if not title:
                continue
            if title.lower() in TEMPLATE_PROMPTS_EXACT:
                add(findings, number, "placeholder-residue",
                    f'de kop is nog het sjabloonwoord "{title}"', "critical")
            elif not is_caps(title):
                add(findings, number, "title-caps",
                    f'de kop staat niet in kapitalen: "{title[:60]}"')

        if layout not in INTENTIONALLY_BLANK_LAYOUTS and not has_text and not any(
            getattr(s, "has_chart", False) or getattr(s, "has_table", False)
            or getattr(s, "shape_type", None) == 13
            for s in slide.shapes
        ):
            add(findings, number, "empty-slide",
                "geen tekst, geen grafiek, geen tabel, geen beeld", "critical")

    # Een aandachtstrekker op elke slide is geen aandacht meer. Ten hoogste één slide op
    # drie draagt grote letter; op de andere slides is de drager gewicht en kleur, of de
    # compositie zelf. Eén drager mag altijd, ook in een deck van twee contentslides.
    ruimte = max(1, round(contentslides / 3)) if contentslides else 1
    if len(slides_met_drager) > ruimte:
        add(findings, 0, "drager-te-vaak",
            f"{len(slides_met_drager)} van de {contentslides} contentslides dragen letter "
            f"van {DRAGER_VLOER:.0f}pt of groter (slides {slides_met_drager}), en er is "
            f"ruimte voor {ruimte}. Op de andere slides is de drager gewicht en kleur — "
            "18pt SemiBold in de hue van zijn categorie — of de compositie zelf.")

    counts = Counter(f["severity"] for f in findings)
    return {
        "deck": str(deck),
        "slides": len(presentation.slides._sldIdLst),
        "charts": charts,
        "tables": tables,
        "fonts": dict(fonts.most_common()),
        "contentslides": contentslides,
        "dragers": slides_met_drager,
        "findings": findings,
        "counts": {"critical": counts["critical"], "warn": counts["warn"]},
        "verdict": "blocked" if counts["critical"] else "clean" if not findings else "warn",
    }


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck")
    parser.add_argument("--strict", action="store_true",
                        help="ook een waarschuwing geeft exit 1")
    args = parser.parse_args()

    result = analyse(Path(args.deck))
    emit(result)

    if result["counts"]["critical"] or (args.strict and result["counts"]["warn"]):
        sys.exit(1)


if __name__ == "__main__":
    main()
