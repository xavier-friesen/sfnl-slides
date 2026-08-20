"""Add a house-style table to a packed deck.

Usage:
    python add_table.py <deck.pptx> --slide 4 --box 0.48,2.4,12.52,3.2 --data tabel.json
    python add_table.py <deck.pptx> --slide 4 --box ... --data tabel.json --zebra
    python add_table.py <deck.pptx> --slide 4 --box ... --data tabel.json --no-header-fill

`tabel.json`:

    {
      "columns": [
        {"label": "Post",        "align": "l"},
        {"label": "2024",        "align": "r", "width_in": 2.0},
        {"label": "2025",        "align": "r", "width_in": 2.0}
      ],
      "rows": [
        ["Rentelasten", "1,2", "0,9"],
        ["Opbrengsten", "0,4", "1,6"]
      ]
    }

`align` is "l", "c" or "r" and defaults to "r" for every column after the first, because
that is what a financial series needs: figures line up on their last digit. `width_in`
is optional per column; whatever is left over is divided over the columns without one.

WHEN TO USE A TABLE. A financial series — three or more periods against two or more
quantities — belongs in a table, not in prose and not in bullets. Prose hides the
comparison the reader came for. Below that, a sentence or a KPI card is usually better:
a two-by-two table is a table pretending to carry weight.

Style, all of it inherited from the theme and none of it configurable by accident:
header row navy (schemeClr accent6) with white Montserrat SemiBold, body Lato Light
12pt, figures right-aligned, `--zebra` for a light navy tint on alternate rows. Only
schemeClr is written, never an srgbClr, so the table follows the theme. Autofit is off on
every cell: PowerPoint shrinking one cell's text makes the row look broken.

Celtekst lijnt VERTICAAL BOVEN uit (`anchor="t"`). Verticaal gecentreerd — de default van
python-pptx — laat de basislijnen binnen één rij verspringen zodra één cel twee regels
draagt: de buurcellen zakken een halve regel mee, en dan is de rij niet meer te volgen.

The built-in PowerPoint table style is switched off (firstRow and bandRow banding), so
what you see here is exactly what the deck shows, rather than the theme's blue default
fighting our fills.

Randen: alleen horizontale scheidingslijnen BINNEN de tabel — geen buitenrand, geen
verticale lijnen. Dat is wat `brand.md` voorschrijft; het volledig navy raster mét
buitenrand dat de render eerder liet zien kwam uit de standaardtabelstijl van het thema,
die overneemt zodra je `tableStyleId` weghaalt en zelf geen randen zet.

Hoogte: de rijhoogtes worden over de gevraagde `--box`-hoogte verdeeld. PowerPoint
negeert de hoogte van het frame en behandelt de rijhoogte als minimum, dus zonder die
verdeling rendert een tabel van drie rijen ~1.3 in hoog terwijl je 3.70 in vroeg — met
twee inch wit tot de conclusieband. De JSON meldt `height_rendered_in` en
`height_reached`.

Run this after the text slides are built and packed, like add_chart.py. Output is JSON.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from _deck import emit

HEADER_FONT = "Montserrat SemiBold"
BODY_FONT = "Lato Light"

HEADER_PT = 12
BODY_PT = 12

# Navy. accent6 is 201B5C in this theme, the same value as dk2, and resolves without
# depending on the slide's colour map.
HEADER_FILL = MSO_THEME_COLOR.ACCENT_6
HEADER_TEXT = MSO_THEME_COLOR.BACKGROUND_1

BODY_TEXT = MSO_THEME_COLOR.TEXT_1

# Light navy tint for the zebra rows: schemeClr accent6 with lumMod/lumOff, which
# python-pptx writes from `brightness`. 0.92 is light enough to read black text on.
ZEBRA_BRIGHTNESS = 0.92

ALIGNMENTS = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}

CELL_MARGIN_IN = 0.08
# Cellen lijnen boven uit (zie `style_cell`), dus de bovenmarge is de optische ruimte
# tussen de scheidingslijn en de eerste regel. 0.06 in houdt de rij luchtig zonder de
# minimale rijhoogte van 0.32 in te overschrijden.
CELL_TOP_MARGIN_IN = 0.06
ROW_HEIGHT_IN = 0.32
HEADER_HEIGHT_IN = 0.36


def parse_columns(columns: list) -> list[dict]:
    """Normalise the column specs and fill in the default alignment."""
    parsed = []
    for index, column in enumerate(columns):
        if isinstance(column, str):
            column = {"label": column}
        if not isinstance(column, dict) or "label" not in column:
            raise SystemExit(
                f"kolom {index} mist een label — geef {{'label': ..., 'align': 'l|c|r'}}"
            )
        # First column is the row label, the rest are figures.
        align = str(column.get("align", "l" if index == 0 else "r")).lower()
        if align not in ALIGNMENTS:
            raise SystemExit(
                f"kolom {index} heeft align={align!r}; kies l, c of r"
            )
        parsed.append(
            {
                "label": str(column["label"]),
                "align": align,
                "width_in": column.get("width_in"),
            }
        )
    return parsed


def column_widths(columns: list[dict], total_in: float) -> list[float]:
    """Explicit widths where given, the remainder split over the rest."""
    fixed = sum(c["width_in"] for c in columns if c["width_in"])
    free = [c for c in columns if not c["width_in"]]
    if fixed > total_in + 1e-6:
        raise SystemExit(
            f"de opgegeven kolombreedtes zijn samen {fixed:.2f} in, breder dan de "
            f"tabel ({total_in:.2f} in)"
        )
    share = (total_in - fixed) / len(free) if free else 0.0
    if free and share <= 0:
        raise SystemExit(
            "geen ruimte over voor de kolommen zonder width_in — geef een bredere box"
        )
    return [c["width_in"] if c["width_in"] else share for c in columns]


def style_cell(cell, *, bold: bool, font: str, size_pt: int, colour, align) -> None:
    frame = cell.text_frame
    frame.word_wrap = True
    # Off, always. A shrunk cell makes the whole row read as a mistake.
    frame.auto_size = MSO_AUTO_SIZE.NONE
    # Boven uitlijnen, niet centreren. Met verticaal gecentreerde cellen verspringen de
    # basislijnen binnen één rij zodra één cel twee regels draagt: de buurcellen zakken
    # een halve regel mee en het oog kan de rij niet meer volgen. Met `anchor="t"` staat
    # elke eerste regel van een rij op dezelfde hoogte, en dat is precies waarvoor een
    # financiële tabel er staat.
    cell.vertical_anchor = MSO_ANCHOR.TOP
    cell.margin_left = Inches(CELL_MARGIN_IN)
    cell.margin_right = Inches(CELL_MARGIN_IN)
    cell.margin_top = Inches(CELL_TOP_MARGIN_IN)
    cell.margin_bottom = Inches(0.03)

    for paragraph in frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.italic = False
            run.font.color.theme_color = colour


def fill_cell(cell, colour, brightness: float | None = None) -> None:
    cell.fill.solid()
    cell.fill.fore_color.theme_color = colour
    if brightness is not None:
        cell.fill.fore_color.brightness = brightness


def clear_fill(cell) -> None:
    cell.fill.background()


def disable_builtin_banding(table) -> None:
    """Strip PowerPoint's own table styling, so only what we set here shows.

    Two steps. The flags switch off the first-row/banding PARTS of the style; removing
    `<a:tableStyleId>` drops the rest — python-pptx writes the "Medium Style 2 - Accent
    1" id by default, which paints a blue header and its own banding on top of ours.

    Wat er dan overblijft is NIET niets: PowerPoint valt terug op de standaardtabelstijl
    van het thema, en in dit sjabloon is dat een volledig navy raster mét buitenrand.
    `brand.md` schrijft juist voor: rechte hoeken, géén buitenrand om de tabel heen. Dus
    zetten we de randen expliciet (zie `style_borders`) in plaats van erop te vertrouwen
    dat "niets schrijven" ook niets oplevert.
    """
    table.first_row = False
    table.horz_banding = False
    table.first_col = False
    table.last_row = False
    table.last_col = False

    tbl_pr = table._tbl.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}tblPr"
    )
    if tbl_pr is not None:
        style_id = tbl_pr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"
        )
        if style_id is not None:
            tbl_pr.remove(style_id)


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _qn(tag: str) -> str:
    return f"{{{A_NS}}}{tag}"


def style_borders(table, rows: int, columns: int) -> None:
    """Randen expliciet: alleen horizontale scheidingslijnen binnen de tabel.

    Per cel staat de lijnkeuze in `<a:tcPr>` als `lnL/lnR/lnT/lnB`, en de volgorde is
    voorgeschreven (L, R, T, B, TlToBr, BlToTr). `<a:noFill/>` is uit, een dunne navy
    lijn is aan.

    Regime, conform `brand.md`:

    * geen buitenrand — de bovenrand van rij 0 en de onderrand van de laatste rij zijn uit
    * geen verticale lijnen — de rechteruitlijning van de cijfers doet dat werk al
    * wel een lijn ONDER de headerrij en tussen de body-rijen: die helpt het oog een rij
      over vier kolommen volgen, en dat was de reden om het raster ooit te laten staan
    """
    from lxml import etree

    hairline_emu = 6350  # 0.5pt

    def set_line(tc, side: str, on: bool, white: bool = False) -> None:
        tc_pr = tc.find(_qn("tcPr"))
        if tc_pr is None:
            tc_pr = etree.SubElement(tc, _qn("tcPr"))
        tag = _qn(side)
        existing = tc_pr.find(tag)
        if existing is not None:
            tc_pr.remove(existing)
        line = etree.Element(tag)
        if on:
            line.set("w", str(hairline_emu))
            line.set("cap", "flat")
            line.set("cmpd", "sng")
            line.set("algn", "ctr")
            fill = etree.SubElement(line, _qn("solidFill"))
            scheme = etree.SubElement(fill, _qn("schemeClr"))
            # Op de navy headerrij moet de scheidingslijn wit zijn, anders zie je hem
            # niet; in de body is hij navy in een lichte tint.
            scheme.set("val", "lt2" if white else "accent6")
            if not white:
                etree.SubElement(scheme, _qn("lumMod")).set("val", "40000")
                etree.SubElement(scheme, _qn("lumOff")).set("val", "60000")
        else:
            etree.SubElement(line, _qn("noFill"))

        # tcPr is een sequence: lnL, lnR, lnT, lnB, lnTlToBr, lnBlToTr, cell3D, ... en
        # daarna de vulling. Op de juiste plaats inserten, niet aanhangen.
        order = ["lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr"]
        position = 0
        for name in order[: order.index(side)]:
            if tc_pr.find(_qn(name)) is not None:
                position += 1
        tc_pr.insert(position, line)

    for row_index in range(rows):
        for col_index in range(columns):
            tc = table.cell(row_index, col_index)._tc
            last_row = row_index == rows - 1
            set_line(tc, "lnL", False)
            set_line(tc, "lnR", False)
            set_line(tc, "lnT", False)
            set_line(tc, "lnB", not last_row, white=row_index == 0)


def row_heights(total_in: float, rows: int) -> tuple[float, float, bool]:
    """(headerhoogte, rijhoogte, gehaald) voor een tabel van `total_in` hoog.

    PowerPoint behandelt de rijhoogte als een MINIMUM en negeert de hoogte van de
    `graphicFrame`: met vaste rijhoogtes van 0.32 in rendert een tabel van drie rijen
    ~1.3 in hoog, ook als je 3.70 in vraagt. Dan staat er tot de conclusieband op 6.93 in
    ruim twee inch wit, en geen enkele check ziet het: een meting leest de opgegeven
    framehoogte en niet de gerenderde. Alleen de render laat dat gat zien.

    Dus verdelen we de gevraagde hoogte over de rijen. Past de inhoud niet in de
    gevraagde hoogte (minder dan de minimummaten), dan is de gevraagde hoogte
    onhaalbaar: dat meldt de JSON, en de tabel wordt zo compact als toegestaan.
    """
    minimum = HEADER_HEIGHT_IN + ROW_HEIGHT_IN * rows
    if total_in <= minimum + 0.01:
        return HEADER_HEIGHT_IN, ROW_HEIGHT_IN, total_in >= minimum - 0.01
    share = (total_in - HEADER_HEIGHT_IN) / rows
    return HEADER_HEIGHT_IN, share, True


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck", type=Path)
    parser.add_argument("--slide", type=int, required=True, help="1-based slide number")
    parser.add_argument("--box", required=True, help="x,y,w,h in inches")
    parser.add_argument("--data", type=Path, required=True, help="table JSON")
    parser.add_argument(
        "--zebra", action="store_true", help="light tint on alternate body rows"
    )
    parser.add_argument(
        "--no-header-fill",
        action="store_true",
        help="header row without the navy fill (navy text on white instead)",
    )
    parser.add_argument("--out", type=Path, help="write here instead of in place")
    args = parser.parse_args()

    # utf-8-sig: op Windows schrijft PowerShells Out-File een BOM, en die is geen reden
    # om een verder geldige tabel.json te weigeren.
    payload = json.loads(args.data.read_text(encoding="utf-8-sig"))
    columns = parse_columns(payload["columns"])
    rows = payload["rows"]
    if not rows:
        raise SystemExit("de tabel heeft geen rijen")
    for index, row in enumerate(rows):
        if len(row) != len(columns):
            raise SystemExit(
                f"rij {index} heeft {len(row)} cellen, de tabel heeft "
                f"{len(columns)} kolommen"
            )

    try:
        x, y, width, height = (float(v) for v in args.box.split(","))
    except ValueError:
        raise SystemExit("--box wil x,y,w,h in inches, bijvoorbeeld 0.48,2.4,12.52,3.2") from None

    widths = column_widths(columns, width)

    presentation = Presentation(str(args.deck))
    if args.slide < 1 or args.slide > len(presentation.slides):
        raise SystemExit(
            f"slide {args.slide} bestaat niet (deck heeft {len(presentation.slides)})"
        )
    slide = presentation.slides[args.slide - 1]

    shape = slide.shapes.add_table(
        len(rows) + 1,
        len(columns),
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    table = shape.table
    disable_builtin_banding(table)

    for index, column_width in enumerate(widths):
        table.columns[index].width = Inches(column_width)

    # `table.rows` is not sliceable in python-pptx, so index it.
    header_in, row_in, height_reached = row_heights(height, len(rows))
    table.rows[0].height = Inches(header_in)
    for row_index in range(1, len(rows) + 1):
        table.rows[row_index].height = Inches(row_in)

    # Header
    for index, column in enumerate(columns):
        cell = table.cell(0, index)
        cell.text = column["label"]
        if args.no_header_fill:
            clear_fill(cell)
            colour = HEADER_FILL
        else:
            fill_cell(cell, HEADER_FILL)
            colour = HEADER_TEXT
        style_cell(
            cell,
            bold=True,
            font=HEADER_FONT,
            size_pt=HEADER_PT,
            colour=colour,
            align=ALIGNMENTS[column["align"]],
        )

    # Body
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index + 1, col_index)
            cell.text = "" if value is None else str(value)
            if args.zebra and row_index % 2 == 1:
                fill_cell(cell, HEADER_FILL, ZEBRA_BRIGHTNESS)
            else:
                clear_fill(cell)
            style_cell(
                cell,
                bold=False,
                font=BODY_FONT,
                size_pt=BODY_PT,
                colour=BODY_TEXT,
                align=ALIGNMENTS[columns[col_index]["align"]],
            )

    style_borders(table, len(rows) + 1, len(columns))

    out = args.out or args.deck
    presentation.save(str(out))
    rendered_height = header_in + row_in * len(rows)
    emit(
        {
            "deck": str(out),
            "slide": args.slide,
            "rows": len(rows),
            "columns": [c["label"] for c in columns],
            "column_widths_in": [round(w, 2) for w in widths],
            "zebra": bool(args.zebra),
            "header_fill": not args.no_header_fill,
            "box_in": [x, y, width, height],
            # De rijhoogtes zijn over de gevraagde hoogte verdeeld; dit is wat PowerPoint
            # werkelijk gaat tekenen. Wijkt `height_rendered_in` af van de vierde waarde
            # van `box_in`, dan was de gevraagde hoogte lager dan de inhoud toelaat.
            "row_heights_in": {"header": round(header_in, 2), "body": round(row_in, 2)},
            "height_rendered_in": round(rendered_height, 2),
            "height_reached": bool(height_reached),
            "borders": "alleen horizontale lijnen binnen de tabel, geen buitenrand",
            "note": (
                None
                if height_reached
                else f"gevraagde hoogte {height:.2f} in is lager dan de minimale "
                f"{rendered_height:.2f} in voor {len(rows)} rijen — de tabel wordt "
                f"{rendered_height:.2f} in. Minder rijen, of geef een hogere box."
            ),
        }
    )


if __name__ == "__main__":
    main()
