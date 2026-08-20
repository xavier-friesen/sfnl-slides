"""Add a native, editable PowerPoint chart to a packed deck.

Usage:
    python add_chart.py <deck.pptx> --slide 4 --data chart.json
    python add_chart.py <deck.pptx> --slide 4 --data chart.json --type line --placeholder 19
    python add_chart.py <deck.pptx> --slide 4 --data chart.json --box 0.48,2.4,6.0,4.0

`chart.json`:

    {
      "categories": ["2023", "2024", "2025"],
      "series": {"Kosten": [1.2, 1.1, 0.9], "Opbrengsten": [0.4, 0.8, 1.6]},
      "number_format": "#,##0.0"
    }

Native charts, not images: a colleague can click the chart in PowerPoint and edit the
numbers. python-pptx writes the chart part and the embedded worksheet, so no XML has
to be hand-built.

Series colours come from the theme, so they are the SFNL colours and stay right if the
theme changes. The default order is orange, royal, sky, emerald, grapefruit. Override
per series with `--series-colors`, which takes the brand names:

    navy, oranje/orange, grapefruit, royal, sky, emerald   (or accent1..accent6, dk1, dk2)

    --series-colors navy,oranje        realisatie navy, prognose oranje

That realisatie-navy / prognose-oranje pairing is the house coding for a chart that
shows realisation against a forecast; use it whenever the chart has that meaning.

`--highlight <categorie>` is the focus colour from `adviesvorm.md` §2: on a
single-series bar chart it puts that one category in orange and every other bar in
navy, so the chart points at its own message. It refuses multi-series charts and
non-bar types instead of colouring something ambiguous:

    --highlight 2025                   de staaf van 2025 oranje, de rest navy

Data labels are on by default when the chart has 8 points or fewer and off above that —
`--labels` forces them on, `--no-labels` forces them off. Chart text is Lato Light.
Gridlines, chart title and shadows stay off: the slide title already says what the
chart means, and the unit belongs in a footnote line on the slide, not in the chart.

The chart is placed in an object placeholder when you name one, otherwise in an
explicit box in inches (`x,y,w,h`), otherwise across the content zone.

De hoogte wordt altijd ingekort tot de onderrand van de contentzone (6.93 in). Dat is
geen smaakkwestie: neem je de doosmaat van een layout-21/22-placeholder over, dan loopt
die zelf tot 6.96 in en staat de grafiek 0.03 in onder de zone, tegen de geërfde chrome
aan. De JSON meldt `clamped_to_zone` en de werkelijke `box_in`; er is geen script dat dat
narekent — de zone-overtreding zie je op de render, en `deck-visual-reviewer` noemt hem
onder overloop.

`c:axId` and `c:crossAx`: python-pptx writes these axis identifiers as NEGATIVE 32-bit
values. This script never touches them and neither should you. "Repairing" them to
positive numbers produces a deck that PowerPoint refuses to open — it has happened
here before. The validator whitelists exactly these two attributes for that reason;
every other chart error is real.

A pie or doughnut plots one series only (python-pptx silently keeps the first and drops
the rest), so this script refuses more than one series for those types instead of
losing data quietly.

Run this after the text slides are built and packed. Output is compact JSON.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor  # noqa: F401  (kept for explicit-colour overrides)
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt

from _deck import emit

CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar_h": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_h_stacked": XL_CHART_TYPE.BAR_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
}
# No "scatter": argparse used to advertise XY_SCATTER, but main() always builds a
# CategoryChartData and python-pptx routes XY_SCATTER to _XyChartXmlWriter, which reads
# series.xValues/.yValues — attributes CategorySeriesData does not have. So --type
# scatter raised inside python-pptx instead of producing a chart. A scatter plot needs
# XyChartData and a different JSON shape; until that is asked for, it is not offered.

# Types whose series mark IS the line: colour the line, never the fill. Styling these
# the fill way wrote <a:ln><a:noFill/></a:ln> onto every series, which deletes the line
# itself and renders an empty plot area that every QA script passes.
LINE_FAMILY = {"line", "line_markers"}

# Types that plot a single series; python-pptx keeps chart_data[0] and drops the rest.
SINGLE_SERIES_ONLY = {"pie", "doughnut"}

# Orange first: it is the SFNL accent. Then royal, sky, emerald, grapefruit.
SERIES_COLOURS = [
    MSO_THEME_COLOR.ACCENT_1,
    MSO_THEME_COLOR.ACCENT_3,
    MSO_THEME_COLOR.ACCENT_4,
    MSO_THEME_COLOR.ACCENT_5,
    MSO_THEME_COLOR.ACCENT_2,
]

# Brand names for --series-colors. navy is accent6 (201B5C), which the theme defines
# identically to dk2 — accent6 is used because a chart part resolves schemeClr against
# the theme without the slide's colour map, where dk2/tx2 can flip.
COLOUR_NAMES = {
    "oranje": MSO_THEME_COLOR.ACCENT_1,
    "orange": MSO_THEME_COLOR.ACCENT_1,
    "grapefruit": MSO_THEME_COLOR.ACCENT_2,
    "royal": MSO_THEME_COLOR.ACCENT_3,
    "sky": MSO_THEME_COLOR.ACCENT_4,
    "emerald": MSO_THEME_COLOR.ACCENT_5,
    "navy": MSO_THEME_COLOR.ACCENT_6,
    "accent1": MSO_THEME_COLOR.ACCENT_1,
    "accent2": MSO_THEME_COLOR.ACCENT_2,
    "accent3": MSO_THEME_COLOR.ACCENT_3,
    "accent4": MSO_THEME_COLOR.ACCENT_4,
    "accent5": MSO_THEME_COLOR.ACCENT_5,
    "accent6": MSO_THEME_COLOR.ACCENT_6,
    "dk1": MSO_THEME_COLOR.DARK_1,
    "dk2": MSO_THEME_COLOR.DARK_2,
}

# Data labels help up to this many plotted points; beyond it they turn into noise.
LABEL_POINT_LIMIT = 8

# --highlight colours per-point fills, which only means something on a bar mark.
HIGHLIGHT_TYPES = {"bar", "bar_h"}


def apply_highlight(chart, categories: list, category: str) -> int:
    """Focus colour: the named category orange, every other point navy."""
    if category not in categories:
        raise SystemExit(
            f"--highlight {category!r} staat niet in de categorieën: "
            + ", ".join(str(c) for c in categories)
        )
    focus = categories.index(category)
    series = chart.plots[0].series[0]
    for index, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.theme_color = (
            MSO_THEME_COLOR.ACCENT_1 if index == focus else MSO_THEME_COLOR.ACCENT_6
        )
        point.format.line.fill.background()
    return focus

BODY_FONT = "Lato Light"
CONTENT_ZONE = (0.48, 1.93, 12.52, 5.00)
LINE_WIDTH_PT = 2.25

# De onderrand van de contentzone. `brand.md` en `xml-editing.md` noemen 6.93 in al de
# ondergrens voor een graphicFrame; hier wordt hij ook afgedwongen. De
# kolomplaceholders van layout 21/22 lopen zelf tot 6.96 in — neem je hun doosmaat
# letterlijk over, dan staat de grafiek 0.03 in te laag, onder de geërfde chrome.
ZONE_BOTTOM_IN = 6.93
EMU_PER_INCH = 914400


def clamp_to_zone(top: int, height: int) -> tuple[int, bool]:
    """Hoogte inkorten tot de contentzone. Returns (hoogte, ingekort)."""
    limit = int(ZONE_BOTTOM_IN * EMU_PER_INCH)
    if top is None or height is None or top + height <= limit:
        return height, False
    return max(limit - top, 1), True


def parse_series_colours(spec: str | None) -> list | None:
    """Turn `navy,oranje` into theme colours, or fail with the list of valid names."""
    if not spec:
        return None
    wanted = []
    for raw in spec.split(","):
        name = raw.strip().lower()
        if not name:
            continue
        if name not in COLOUR_NAMES:
            raise SystemExit(
                f"onbekende kleur {raw.strip()!r} — kies uit: "
                + ", ".join(sorted(COLOUR_NAMES))
            )
        wanted.append(COLOUR_NAMES[name])
    return wanted or None


def style_text(font, size_pt: int = 12) -> None:
    font.name = BODY_FONT
    font.size = Pt(size_pt)
    font.bold = False
    font.italic = False
    font.color.theme_color = MSO_THEME_COLOR.TEXT_1


def style_chart(
    chart,
    labels: bool,
    legend: str | None,
    number_format: str | None,
    chart_type: str = "bar",
    series_colours: list | None = None,
) -> None:
    chart.has_title = False
    style_text(chart.font)

    single_series = len(chart.plots[0].series) == 1
    show_legend = legend != "none" and not single_series
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = {
            "bottom": XL_LEGEND_POSITION.BOTTOM,
            "right": XL_LEGEND_POSITION.RIGHT,
            "top": XL_LEGEND_POSITION.TOP,
        }.get(legend or "bottom", XL_LEGEND_POSITION.BOTTOM)
        chart.legend.include_in_layout = False
        style_text(chart.legend.font)

    palette = series_colours or SERIES_COLOURS
    draws_a_line = chart_type in LINE_FAMILY

    for index, series in enumerate(chart.plots[0].series):
        colour = palette[index % len(palette)]
        try:
            if draws_a_line:
                # series.format.line IS the mark here. Colour it and leave the fill
                # alone — a line series has no meaningful area fill.
                series.format.line.color.theme_color = colour
                series.format.line.width = Pt(LINE_WIDTH_PT)
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.theme_color = colour
                # For a bar/column/pie/area series this is only the outline, so
                # removing it is the clean look we want.
                series.format.line.fill.background()
        except (NotImplementedError, ValueError):
            pass

    plot = chart.plots[0]
    plot.has_data_labels = labels
    if labels:
        plot.data_labels.show_value = True
        style_text(plot.data_labels.font, 11)
        if number_format:
            plot.data_labels.number_format = number_format
            plot.data_labels.number_format_is_linked = False

    for axis_name in ("category_axis", "value_axis"):
        try:
            axis = getattr(chart, axis_name)
        except (ValueError, NotImplementedError):
            continue
        axis.has_major_gridlines = False
        axis.has_minor_gridlines = False
        axis.major_tick_mark = XL_TICK_MARK.NONE
        axis.minor_tick_mark = XL_TICK_MARK.NONE
        if axis.has_title:
            axis.has_title = False
        style_text(axis.tick_labels.font)
        if axis_name == "value_axis" and number_format:
            axis.tick_labels.number_format = number_format
            axis.tick_labels.number_format_is_linked = False


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck", type=Path)
    parser.add_argument("--slide", type=int, required=True, help="1-based slide number")
    parser.add_argument("--data", type=Path, required=True, help="chart JSON")
    parser.add_argument("--type", default="bar", choices=sorted(CHART_TYPES))
    parser.add_argument("--placeholder", type=int, help="object placeholder idx to fill")
    parser.add_argument("--box", help="x,y,w,h in inches")
    parser.add_argument(
        "--labels",
        action="store_true",
        help=f"force data labels on (default: on at <= {LABEL_POINT_LIMIT} points)",
    )
    parser.add_argument(
        "--no-labels", action="store_true", help="force data labels off"
    )
    parser.add_argument(
        "--series-colors",
        dest="series_colors",
        help="per series, e.g. navy,oranje (realisatie, prognose)",
    )
    parser.add_argument(
        "--highlight",
        help="focuskleur: deze categorie oranje, de rest navy (één reeks, bar/bar_h)",
    )
    parser.add_argument("--legend", default="bottom", choices=["bottom", "right", "top", "none"])
    parser.add_argument("--out", type=Path, help="write here instead of in place")
    args = parser.parse_args()

    # utf-8-sig: op Windows schrijft PowerShells Out-File een BOM, en die is geen reden
    # om een verder geldige chart.json te weigeren.
    payload = json.loads(args.data.read_text(encoding="utf-8-sig"))
    categories = payload["categories"]
    series = payload["series"]
    number_format = payload.get("number_format")

    if args.type in SINGLE_SERIES_ONLY and len(series) > 1:
        raise SystemExit(
            f"een {args.type}-diagram plot één reeks; je geeft er {len(series)} "
            f"({', '.join(series)}). Kies bar of splits de slide."
        )

    if args.highlight is not None:
        if args.type not in HIGHLIGHT_TYPES:
            raise SystemExit(
                f"--highlight kleurt staven en werkt niet op {args.type} — "
                f"kies uit: {', '.join(sorted(HIGHLIGHT_TYPES))}"
            )
        if len(series) > 1:
            raise SystemExit(
                f"--highlight is de focuskleur op één reeks; je geeft er {len(series)}. "
                "Gebruik --series-colors om reeksen te coderen."
            )
        if args.series_colors:
            raise SystemExit("--highlight en --series-colors sluiten elkaar uit")

    series_colours = parse_series_colours(args.series_colors)
    if series_colours and len(series_colours) < len(series):
        raise SystemExit(
            f"--series-colors geeft {len(series_colours)} kleuren voor "
            f"{len(series)} reeksen"
        )

    points = len(categories) * max(len(series), 1)
    if args.no_labels:
        labels = False
    elif args.labels:
        labels = True
    else:
        labels = points <= LABEL_POINT_LIMIT

    data = CategoryChartData()
    data.categories = categories
    for name, values in series.items():
        data.add_series(name, values, number_format)

    presentation = Presentation(str(args.deck))
    if args.slide < 1 or args.slide > len(presentation.slides):
        raise SystemExit(f"slide {args.slide} bestaat niet (deck heeft {len(presentation.slides)})")
    slide = presentation.slides[args.slide - 1]

    placed_in = None
    clamped = False
    if args.placeholder is not None:
        target = next(
            (p for p in slide.placeholders if p.placeholder_format.idx == args.placeholder),
            None,
        )
        if target is None:
            raise SystemExit(
                f"placeholder {args.placeholder} staat niet op slide {args.slide}"
            )
        if hasattr(target, "insert_chart"):
            graphic_frame = target.insert_chart(CHART_TYPES[args.type], data)
            placed_in = f"placeholder {args.placeholder}"
        else:
            # Placeholders die add_slide.py uit een layout overneemt komen als gewone
            # SlidePlaceholder terug, zonder insert_chart. De grafiek gaat dan op de
            # geometrie van die placeholder staan; de placeholder zelf verdwijnt,
            # zodat er geen leeg tekstvak achterblijft.
            left, top, width, height = (
                target.left, target.top, target.width, target.height
            )
            if None in (left, top, width, height):
                raise SystemExit(
                    f"placeholder {args.placeholder} op slide {args.slide} heeft geen "
                    "geometrie die uit de layout te herleiden is — geef --box "
                    "x,y,w,h in inches"
                )
            height, clamped = clamp_to_zone(top, height)
            target._element.getparent().remove(target._element)
            graphic_frame = slide.shapes.add_chart(
                CHART_TYPES[args.type], left, top, width, height, data
            )
            placed_in = f"placeholder {args.placeholder} (geometrie overgenomen)"
    else:
        if args.box:
            x, y, width, height = (float(v) for v in args.box.split(","))
        else:
            x, y, width, height = CONTENT_ZONE
        emu_height, clamped = clamp_to_zone(int(Inches(y)), int(Inches(height)))
        graphic_frame = slide.shapes.add_chart(
            CHART_TYPES[args.type], Inches(x), Inches(y), Inches(width), emu_height, data
        )
        placed_in = f"box {x},{y},{width},{height}"

    style_chart(
        graphic_frame.chart,
        labels,
        args.legend,
        number_format,
        chart_type=args.type,
        series_colours=series_colours,
    )
    if args.highlight is not None:
        apply_highlight(graphic_frame.chart, categories, args.highlight)

    out = args.out or args.deck
    presentation.save(str(out))
    emit(
        {
            "deck": str(out),
            "slide": args.slide,
            "type": args.type,
            "series": list(series),
            "categories": len(categories),
            "labels": labels,
            "highlight": args.highlight,
            "placed_in": placed_in,
            "box_in": [
                round(graphic_frame.left / EMU_PER_INCH, 3),
                round(graphic_frame.top / EMU_PER_INCH, 3),
                round(graphic_frame.width / EMU_PER_INCH, 3),
                round(graphic_frame.height / EMU_PER_INCH, 3),
            ],
            "clamped_to_zone": clamped,
            "note": (
                f"hoogte ingekort tot de contentzone (onderrand {ZONE_BOTTOM_IN} in) — "
                "de overgenomen doos liep daaronder"
                if clamped
                else None
            ),
        }
    )


if __name__ == "__main__":
    main()
