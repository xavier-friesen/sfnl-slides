"""Shared reading helpers for the QA scripts.

The point of this module is `resolve_run`: in a layout-first deck almost nothing is set
on the run itself. Size, font and colour come from the layout placeholder, the master
placeholder or the master text styles. QA that only looks at `run.font.size` sees None
everywhere and finds nothing.

The resolution order follows what PowerPoint does:

    run rPr
    → paragraph pPr/defRPr
    → the shape's own lstStyle for that level
    → the matching placeholder in the layout
    → the matching placeholder in the master
    → the master's txStyles (titleStyle / bodyStyle / otherStyle)
    → presentation defaultTextStyle
    → 18pt Calibri Light

Dit is óók de plek van het LAYOUTBELEID — `LAYOUT_POLICY`, op layoutNAAM — en van de
sjabloonfeiten die uit de gegenereerde sidecar `reference/layouts.json` komen. Zie het
blok met dat opschrift hieronder: daar staat waarom het beleid op naam keyt en niet op
nummer, en waar nummers wél correct zijn.
"""

from __future__ import annotations

import glob
import json
import os
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS = {"a": A, "p": P, "r": R}

EMU_PER_INCH = 914400

FALLBACK_SIZE_PT = 18.0
NEUTRAL_HEX = {"FFFFFF", "FEFFFF", "000000"}

# `Lato Semibold` staat erbij sinds de aanhef binnen een doorlopende regel Lato is en niet
# Montserrat (`vormentaal.md` §9): één familie per regel, en dan hoort er een echt zwaarder
# Lato-gewicht te bestaan in plaats van `b="1"` op Lato Light.
BRAND_FONTS = {"Gotham Bold", "Montserrat", "Montserrat Light", "Montserrat SemiBold",
               "Montserrat Bold", "Lato Light", "Lato", "Lato Semibold", "Lato Bold"}

SCRIPT_DIR = Path(__file__).resolve().parent


# ---------------------------------------------------------------------------
# Layoutbeleid (FIXSPEC addendum B11, bevestigd door de eigenaar 2026-07-29; op NAAM
# gezet in W2b).
#
# HET BELEID KEYT OP LAYOUTNAAM, NIET OP NUMMER. Het nummer komt uit de partname
# `slideLayoutN.xml` en is alleen binnen ÉÉN sjablooninstantie geldig. Het referentiedeck
# `Impact meetplan Welzijn op Recept.pptx` heeft 3 masters en 44 layouts, en onze
# norm-contentslide `'Titel, subtitel'` bestaat daar als layout 20, 23 én 37. Zolang het
# beleid op nummer keyde, werd elke correcte contentslide van dat deck op
# `slideLayout23.xml` gemeld als kritieke `forbidden-layout` — precies de melding die
# `sfnl-ppt-fix/SKILL.md` de agent verbiedt weg te argumenteren. De NAAM is stabiel over
# sjablooninstanties; het nummer is dat niet.
#
# Waar nummers WEL correct zijn: ons eigen bouwpad, waar het sjabloon bekend is.
# `add_slide.py slideLayout19.xml`, `prune_template.py --drop-layouts 18,23,24`, de
# outline en `deck_spec.py` (een spec noemt óns layoutnummer). Daar staan de
# nummer-gesleutelde verzamelingen onderaan dit blok voor, en die worden AFGELEID uit de
# beleidstabel plus de sjabloonfeiten — niet apart onderhouden.
#
# ÉÉN BRON VAN WAARHEID, in twee helften die elkaar niet overlappen:
#
#   1. BELEID — `LAYOUT_POLICY` hieronder, op naam. Verdict, wat er minimaal gevuld moet
#      zijn, de reden bij een verboden layout, en het foto-onderwerp. Dat zijn oordelen;
#      geen script kan ze uit een .potx lezen.
#   2. SJABLOONFEITEN — `reference/layouts.json`, GEGENEREERD door `layout_catalog.py`
#      uit `assets/sfnl-sjabloon.potx`. Nummer, naam, master, tekstplaceholders in
#      leesvolgorde, placeholdermaten, kleurvlak. Dat leest een script wél, en dus hoort
#      het nergens met de hand te staan. `deck_spec.py` hield hier tot W2b een eigen kopie
#      van (`LAYOUT_PLACEHOLDERS`, `LAYOUT_REQUIRED`, `MASTER_ONE_LAYOUTS`,
#      `AGENDA_COLOURS`, `LAYOUT_BOXES`); die leest hij nu hieruit.
#
# De twee worden op NAAM aan elkaar geknoopt. `test_catalog.py` toetst dat de
# meegeleverde sidecar nog byte-identiek regenereert uit het sjabloon, dus de feitenhelft
# kan niet wegdrijven; `test_deck_spec.py` toetst dat elke layout in het sjabloon een
# beleidsregel heeft, dus de beleidshelft ook niet.
# ---------------------------------------------------------------------------

#: Verdict voor een layoutnaam die in `LAYOUT_POLICY` niet voorkomt. Expliciet, want de
#: twee stille alternatieven zijn beide fout: `nooit` beschuldigt een deck van een
#: overtreding die niemand kan nagaan, en `norm` doet alsof een vreemde layout onze
#: contentslide is. `onbekend` zegt wat het is: dit deck komt niet uit ons sjabloon.
UNKNOWN_VERDICT = "onbekend"


@dataclass(frozen=True)
class LayoutRule:
    """Het beleid voor één layoutNAAM.

    `forbidden` is (wat het is, wat je in plaats daarvan doet) en is alleen gezet bij
    verdict `nooit`; `required` is wat er minimaal gevuld moet zijn wil de slide iets
    zeggen (niet álle placeholders — de subtitel is per beleid optioneel); `photo`
    beschrijft wat er op de layoutfoto te zien is.
    """

    verdict: str
    required: frozenset[int] = field(default_factory=frozenset)
    photo: str | None = None
    forbidden: tuple[str, str] | None = None


_TABEL_LOOK = (
    "kolommen met oranje kopbalken en omlijnde vakken (de tabel-look)",
    "componeer de kolommen zelf op layout 19, of maak er een echte tabel van "
    "met add_table.py",
)

#: Het volledige layoutbeleid, op NAAM. Eén regel per layoutnaam in het SFNL-sjabloon.
#: De nummers in het commentaar zijn die van óns sjabloon en staan er om te kunnen
#: nazoeken; ze doen in de code niets.
#:
#: Het foto-onderwerp is HANDGESCHREVEN, na het bekijken van de layoutfoto's in het
#: sjabloon zelf (renders: test/template-analyse/). Geen script kan naar een foto kijken
#: en zeggen wat erop staat. Beschrijf wat er te zien is, niet wat het zou moeten
#: betekenen. Let op dat naamvolgorde en nummervolgorde bij de sectieslides NIET
#: samenvallen: `1_sectieslide_stijl1` is layout 7 en `2_sectieslide_stijl1` is layout 6.
#:
#: De regels van de elf dividers (6 t/m 16) zijn op 30 juli 2026 HERSCHREVEN op elf verse
#: PowerPoint-renders, omdat ze de foto's gunstiger beschreven dan ze zijn en dat de bouwer
#: renderrondes kostte (W6-meting, bevinding 2): layout 9 heette "wijk, ontmoeting, kleine
#: ondernemers" en is twee stilzittende oudere mannen bij een snackbar, wat boven een
#: hoofdstuk over een groeiende onderneming het omgekeerde zegt. Elke regel heeft nu drie
#: delen, in deze volgorde: (1) wat er te zien is, (2) de dominante kleuren waar die een
#: deck binnenkomen, (3) waar de foto NIET bij past. Dat derde deel is het deel dat de
#: keuze goedkoop maakt — een foto uitsluiten kost één regel lezen, een foto proberen kost
#: een bouw plus een render. Geen aanbevelingen, geen "de meest gebruikte opener voor":
#: fotokeuze is layoutkeuze, en een instrument dat verkoopt misleidt.
LAYOUT_POLICY: dict[str, LayoutRule] = {
    # --- covers (1, 4) ---
    "1_Titelslide": LayoutRule(
        "cover",
        required=frozenset({14}),
        photo="Drukke markt in een winkelstraat, mensen tussen de kramen — stad, publiek, "
              "dagelijks leven. Neutrale externe cover.",
    ),
    "7_Titelslide": LayoutRule("cover", required=frozenset({10})),
    # --- eindslides (2, 3) ---
    "5_Titelslide": LayoutRule("eindslide"),
    "6_Titelslide": LayoutRule("eindslide"),
    # --- quote (5) ---
    "Quote": LayoutRule(
        "quote",
        required=frozenset({0}),
        photo="Vrouw met hoofddoek achter een toonbank met brood, oranje vlak eroverheen — "
              "ondernemerschap, werk, dienstverlening.",
    ),
    # --- sectiedividers, stijl 1 met foto (6-16) ---
    "2_sectieslide_stijl1": LayoutRule(  # 6
        "sectiedivider",
        required=frozenset({14}),
        photo="Jongeman in sportkleding slaat een boksbal in een stadspark; op de bal staat "
              "in rood en blauw 'BOOGIE' geschreven. Groen loof, grijs asfalt, rood op de "
              "bal. Past bij sport, buitenruimte, jongeren. Niet boven een hoofdstuk over "
              "geld, cijfers of organisatie. Deze layout heeft óók een content-placeholder "
              "(idx 11), dus bruikbaar als hoofdstukoverzicht bij de start van een hoofdstuk.",
    ),
    "1_sectieslide_stijl1": LayoutRule(  # 7
        "sectiedivider",
        required=frozenset({14}),
        photo="Glazen bakje vol koperen en zilveren munten waaruit twee groene blaadjes "
              "groeien, tegen een blauwgrijze achtergrond. Koper, groen en blauwgrijs; "
              "brengt geen kleur binnen die buiten het palet valt. Past bij sparen, opbouw, "
              "een fonds. Niet boven een hoofdstuk over mensen, uitvoering of een lopende "
              "onderneming: er staat geen mens op en niets beweegt.",
    ),
    "3_sectieslide_stijl1": LayoutRule(  # 8
        "sectiedivider",
        required=frozenset({14}),
        photo="Jongeman met pet zet een truc op een fiets voor een monumentaal bakstenen "
              "pand, met gras en bomen eromheen. Baksteenrood en groen. Past bij jeugd, "
              "sport, voorzieningen in de wijk. Niet boven een financieel of bestuurlijk "
              "hoofdstuk.",
    ),
    "4_sectieslide_stijl1": LayoutRule(  # 9
        "sectiedivider",
        required=frozenset({14}),
        photo="Twee oudere mannen zitten stil voor een snackbar met de naam 'PITA QUEEN'; de "
              "een op een scootmobiel met een rode en een oranje boodschappentas, de ander "
              "onderuitgezakt op een bankje met een telefoon aan zijn oor. Baksteenrood "
              "trottoir, felrode tas. De twee zitten apart en doen niets: dit leest als "
              "wachten en als een laag inkomen, niet als ontmoeting en niet als "
              "ondernemerschap. Niet boven een hoofdstuk over groei, over een onderneming of "
              "over samenwerking — daar zegt de foto het omgekeerde.",
    ),
    "11_sectieslide_stijl1": LayoutRule(  # 10
        "sectiedivider",
        required=frozenset({14}),
        photo="Hand met een biljet van tien euro en losse munten boven kratten aardbeien op "
              "een markt. Fel rood over de onderste helft plus warm karton; dat rood ligt "
              "naast grapefruit maar is het niet. Past bij dagelijkse uitgaven, huishoudgeld, "
              "besteding. Niet bij een bedrijfsbegroting of een investeringsbedrag: het geld "
              "op deze foto is klein geld.",
    ),
    "6_sectieslide_stijl1": LayoutRule(  # 11
        "sectiedivider",
        required=frozenset({14}),
        photo="Twee handen met rood gelakte nagels typen van bovenaf op een laptop; de "
              "mouwen zijn zachtgeel. Geel en rood, beide buiten het palet, en de laptop is "
              "een merkmodel met het logo in beeld. Past bij administratie, aanvraag, "
              "uitvoering achter de schermen. Niet boven een hoofdstuk over mensen of over "
              "resultaat.",
    ),
    "7_sectieslide_stijl1": LayoutRule(  # 12
        "sectiedivider",
        required=frozenset({14}),
        photo="Twee paar handen sorteren donkere bonen op een turkooizen doek op rode aarde, "
              "met een lichtgroene plant ernaast. Turkoois en roodbruin domineren en staan "
              "naast geen van de zes merkkleuren. Dit is landbouwwerk in een tropisch "
              "ogende omgeving; de jury van ronde 3 wees precies deze foto aan als 'bonen op "
              "rode aarde' boven een provinciaal fonds. Alleen bruikbaar waar het hoofdstuk "
              "letterlijk over landbouw of over werk buiten Nederland gaat.",
    ),
    "8_sectieslide_stijl1": LayoutRule(  # 13
        "sectiedivider",
        required=frozenset({14}),
        photo="Iemand in een beige trenchcoat en rood-witte sneakers zit hoog op een ladder "
              "bij een demonstratie, met een handgeschreven bord 'YOU PAID FOR THIS' ernaast; "
              "herfstbomen en een menigte op de achtergrond. Herfstgeel en oranje, met rood "
              "als accent. De Engelse tekst is leesbaar en spreekt mee: hij verwijt de lezer "
              "iets. Alleen gebruiken als dat verwijt de boodschap is, en nooit in een deck "
              "waarin de lezer de opdrachtgever of de financier is.",
    ),
    "9_sectieslide_stijl1": LayoutRule(  # 14
        "sectiedivider",
        required=frozenset({14}),
        photo="Benen van hardlopers in een wedstrijd op nat asfalt, met de spiegeling in een "
              "plas eronder. Neonkoraal schoenen, limegroene en turkooizen sokken, roze "
              "broek: drie kleuren die in geen SFNL-deck voorkomen en die de slide "
              "domineren. Past bij een wedstrijd of een evenement. Niet als beeld voor een "
              "cohort of een doorlooptijd — het brengt precies de kleur binnen waar de regel "
              "hierboven tegen waarschuwt.",
    ),
    "10_sectieslide_stijl1": LayoutRule(  # 15
        "sectiedivider",
        required=frozenset({14}),
        photo="Brede nieuwe brug over water naar een blok woningen en een bakstenen "
              "appartementengebouw, grijze bewolkte lucht, één fietser en twee voetgangers. "
              "Grijs, baksteenrood en blauwgrijs, alles binnen het palet. De meest neutrale "
              "van de elf: geen leesbare tekst, geen gezicht, geen kleur die eruit springt. "
              "Past bij gebiedsontwikkeling, infrastructuur, een stedelijke opgave, en is de "
              "kandidaat wanneer de uitkomst 'één neutrale foto voor alle hoofdstukken' is.",
    ),
    "12_sectieslide_stijl1": LayoutRule(  # 16
        "sectiedivider",
        required=frozenset({14}),
        photo="Laptop op een tafel met een webanalytics-dashboard in beeld: een blauwe "
              "heatmap, een wereldkaart en een donutdiagram. Blauw en grijs, binnen het "
              "palet. Dit is een stockfoto van software en niet van het werk zelf; de jury "
              "van ronde 3 wees hem af boven een hoofdstuk over rendement. Alleen bruikbaar "
              "waar het hoofdstuk over de monitoring of de rapportage zélf gaat.",
    ),
    # --- het blanco canvas (17) ---
    "Leeg": LayoutRule("speciaal-canvas"),
    # --- contentlayouts (19-22) ---
    "Titel, subtitel": LayoutRule("norm", required=frozenset({0})),                    # 19
    "1_Titel, subtitel, tekst": LayoutRule("toegestaan", required=frozenset({0, 10})),  # 20
    "Titel, subtitel, twee tekstvakken": LayoutRule(                                    # 21
        "toegestaan", required=frozenset({0})
    ),
    "1_Titel, subtitel, twee tekstvakken": LayoutRule(                                  # 22
        "toegestaan", required=frozenset({0})
    ),
    # --- verboden (18, 23, 24) ---
    #
    # Deze drie zijn in W1 uit óns sjabloon verwijderd (`prune_template.py
    # --drop-layouts 18,23,24`), maar het beleid MOET ze blijven kennen: `sfnl-ppt-fix`
    # draait op decks van collega's die op dezelfde merkversie gebouwd zijn en die de
    # layouts wél dragen. Daarom keyt dit op naam en niet op nummer: bij ons zijn 18, 23
    # en 24 gaten, bij een ander deck bezette en soms volkomen geldige layouts.
    "Titel": LayoutRule(                                                                # 18
        "nooit",
        forbidden=(
            "Titel — wit mét titel, zónder de oranje merklijn",
            "gebruik layout 19: dezelfde vrije contentzone, mét de oranje lijn",
        ),
    ),
    "2_Titel, subtitel, twee tekstvakken": LayoutRule("nooit", forbidden=_TABEL_LOOK),   # 23
    "3_Titel, subtitel, twee tekstvakken": LayoutRule("nooit", forbidden=_TABEL_LOOK),   # 24
    # --- agenda's en opsommingen, stijl 2 met kleurvlak (25-30) ---
    **{
        f"{n}_sectieslide_stijl2": LayoutRule(
            "agenda-opsomming", required=frozenset({0, 11})
        )
        for n in range(1, 7)
    },
}

#: Verboden layouts, op NAAM: {naam: (wat het is, wat je in plaats daarvan doet)}. Een
#: slide hierop is een kritieke bevinding (check `forbidden-layout`). Naamsleutel, want
#: dit is precies de check die op een vréémd deck draait.
FORBIDDEN_LAYOUTS: dict[str, tuple[str, str]] = {
    name: rule.forbidden for name, rule in LAYOUT_POLICY.items() if rule.forbidden
}


def _names_with_verdict(*verdicts: str) -> frozenset[str]:
    return frozenset(
        name for name, rule in LAYOUT_POLICY.items() if rule.verdict in verdicts
    )


#: De norm-contentslide. Begin hier, tenzij er een reden is om dat niet te doen.
NORM_LAYOUT_NAME = "Titel, subtitel"

#: De contentlayouts: de norm plus de drie tekstuele varianten. Alleen hierop gelden de
#: titel-/subtitelmodus-checks. De verboden namen staan er NIET in: een modusmelding erop
#: zou de echte bevinding (`forbidden-layout`) verhullen.
CONTENT_LAYOUT_NAMES = _names_with_verdict("norm", "toegestaan")

#: Het bewuste blanco canvas. Heeft GEEN titelplaceholder, dus geen enkel script mag
#: erover klagen dat de titel ontbreekt, en de zone-checks van `qa_fit.py` gelden er niet:
#: er is geen contentzone, er is een heel canvas.
BLANK_CANVAS_NAME = "Leeg"

COVER_LAYOUT_NAMES = _names_with_verdict("cover")
OUTRO_LAYOUT_NAMES = _names_with_verdict("eindslide")
QUOTE_LAYOUT_NAMES = _names_with_verdict("quote")
DIVIDER_LAYOUT_NAMES = _names_with_verdict("sectiedivider")
AGENDA_LAYOUT_NAMES = _names_with_verdict("agenda-opsomming")


def verdict_for_name(name: str | None) -> str:
    """Het verdict van een layoutNAAM, of `onbekend`.

    Dit is de enige plek waar een verdict vandaan komt. `layout_catalog.py` bouwt de
    verdict-kolom hiermee, `inspect_deck.py` rapporteert hem, en `qa_text.py` leidt er de
    `forbidden-layout`-bevinding uit af.
    """
    rule = LAYOUT_POLICY.get(name or "")
    return rule.verdict if rule is not None else UNKNOWN_VERDICT


def forbidden_reason(name: str | None) -> tuple[str, str] | None:
    """(wat het is, wat je in plaats daarvan doet) als deze NAAM verboden is."""
    return FORBIDDEN_LAYOUTS.get(name or "")


def layout_name_of(slide) -> str | None:
    """De naam van de layout achter een python-pptx slide.

    Dit is wat het beleid leest. Zie `layout_number_of` voor waarom niet het nummer.
    """
    try:
        return slide.slide_layout.name
    except AttributeError:
        return None


def layout_number_of(slide) -> int | None:
    """Het layoutNUMMER achter een python-pptx slide, uit zijn partname.

    LET OP: alleen geldig binnen één sjablooninstantie. Gebruik dit om een slide te
    IDENTIFICEREN (in een melding, of om terug te vinden welk part erbij hoort), nooit om
    er beleid aan te hangen — daarvoor is `layout_name_of`.
    """
    try:
        partname = str(slide.slide_layout.part.partname)
    except AttributeError:
        return None
    match = re.search(r"slideLayout(\d+)\.xml", partname)
    return int(match.group(1)) if match else None


# ---------------------------------------------------------------------------
# Sjabloonfeiten — gelezen uit de gegenereerde sidecar `reference/layouts.json`
#
# Alles hieronder is nummer-gesleuteld en beschrijft ÓNS sjabloon. Dat is precies waar
# nummers wél mogen: `add_slide.py` krijgt `slideLayout19.xml` mee, een deck-spec noemt
# ons layoutnummer, en `prune_template.py` noemt 18/23/24. Voor de beoordeling van een
# willekeurig deck gebruik je de naam-API hierboven.
# ---------------------------------------------------------------------------

SIDECAR_PATH = SCRIPT_DIR.parent / "reference" / "layouts.json"

#: De drie layouts die W1 uit het sjabloon haalde. Ze staan niet in de sidecar — die is
#: gegenereerd uit het huidige sjabloon en daar bestaan ze niet meer — terwijl
#: `deck_spec.py` ze wél moet kennen: een spec die layout 23 noemt hoort te horen dat hij
#: de kolommen zelf op 19 componeert, en niet dat layout 23 niet bestaat. Die laatste
#: melding zegt niet wat je moet doen. Naam en tekstplaceholders komen uit het sjabloon
#: van vóór de prune (`git show 8f91fb5^:...`, en `01 SFNL_sjabloon.potx` in de wortel).
REMOVED_LAYOUTS: dict[int, tuple[str, tuple[int, ...]]] = {
    18: ("Titel", (0,)),
    23: ("2_Titel, subtitel, twee tekstvakken", (0, 1, 16, 20, 19, 21)),
    24: ("3_Titel, subtitel, twee tekstvakken", (0, 1, 16, 20, 22, 19, 21, 23)),
}


def _load_layout_facts() -> list[dict]:
    """De layouts uit `reference/layouts.json`, of een lege lijst met een waarschuwing.

    Lege lijst en niet een harde fout, om precies één reden: `layout_catalog.py` is het
    script dat de sidecar GENEREERT en het importeert deze module. Zou het ontbreken van
    de sidecar hier fataal zijn, dan was hij niet meer te regenereren.
    """
    if not SIDECAR_PATH.exists():
        print(
            f"# {SIDECAR_PATH} ontbreekt — de layoutfeiten zijn leeg. Regenereer met "
            "`python scripts/layout_catalog.py`.",
            file=sys.stderr,
        )
        return []
    try:
        data = json.loads(SIDECAR_PATH.read_text(encoding="utf-8"))
    except (OSError, ValueError) as problem:
        print(
            f"# {SIDECAR_PATH} is niet te lezen ({problem}) — regenereer met "
            "`python scripts/layout_catalog.py`.",
            file=sys.stderr,
        )
        return []
    return list(data.get("layouts") or [])


_LAYOUT_FACTS: list[dict] = _load_layout_facts()

#: {19: 'Titel, subtitel'} — het nummer waaronder een naam in ÓNS sjabloon zit, plus de
#: drie verwijderde layouts. Dit is de mapping die bij `assets/sfnl-sjabloon.potx` hoort.
LAYOUT_NAMES: dict[int, str] = {
    **{int(entry["number"]): entry["name"] for entry in _LAYOUT_FACTS},
    **{number: name for number, (name, _phs) in REMOVED_LAYOUTS.items()},
}

#: Tekstplaceholders per layout, in LEESVOLGORDE — dezelfde volgorde als de `--only`-regel
#: in `layouts.md`. Datum- en paginanummerplaceholders staan er niet in: `add_slide.py`
#: zet die nooit op een slide (SKIP_PH_TYPES), dus ze zijn niet te vullen. De layouts 2, 3
#: en 17 hebben daarom een leeg tuple, en dat is correct — daar is niets te vullen.
LAYOUT_PLACEHOLDERS: dict[int, tuple[int, ...]] = {
    **{
        int(entry["number"]): tuple(entry.get("text_placeholders") or ())
        for entry in _LAYOUT_FACTS
    },
    **{number: phs for number, (_name, phs) in REMOVED_LAYOUTS.items()},
}

#: Placeholdermaten per layout, `{layout: {idx: (x, y, breedte, hoogte)}}` in inch, op
#: twee decimalen — dezelfde afronding als de positiekolom in `layouts.md`.
LAYOUT_BOXES: dict[int, dict[int, tuple[float, float, float, float]]] = {
    int(entry["number"]): {
        int(idx): tuple(box) for idx, box in (entry.get("boxes") or {}).items()
    }
    for entry in _LAYOUT_FACTS
}

#: Master 1 draagt het paginanummer als getekend tekstvak in de master, dus
#: `--no-page-number` (showMasterSp="0" op de layout) kan daar. Op master 2 zit het logo
#: óók in de master en weigert `add_slide.py`.
MASTER_ONE_LAYOUTS: frozenset[int] = frozenset(
    int(entry["number"]) for entry in _LAYOUT_FACTS if entry.get("master") == 1
)

#: Het kleurvlak per agenda-layout, uit het sjabloon gelezen. De layouts 25-30 zijn
#: identiek op dat vlak na.
AGENDA_COLOURS: dict[int, str] = {
    int(entry["number"]): entry["colour_block"]
    for entry in _LAYOUT_FACTS
    if entry.get("colour_block") and verdict_for_name(entry["name"]) == "agenda-opsomming"
}


def _numbers_for(names) -> frozenset[int]:
    """De nummers waaronder deze layoutNAMEN in ons sjabloon zitten."""
    return frozenset(number for number, name in LAYOUT_NAMES.items() if name in names)


CONTENT_LAYOUTS = _numbers_for(CONTENT_LAYOUT_NAMES)
COVER_LAYOUTS = _numbers_for(COVER_LAYOUT_NAMES)
OUTRO_LAYOUTS = _numbers_for(OUTRO_LAYOUT_NAMES)
DIVIDER_LAYOUTS = _numbers_for(DIVIDER_LAYOUT_NAMES)
AGENDA_LAYOUTS = _numbers_for(AGENDA_LAYOUT_NAMES)

#: `None` zodra het sjabloon de layout niet meer draagt; in dit sjabloon 5 en 17.
QUOTE_LAYOUT: int | None = next(iter(sorted(_numbers_for(QUOTE_LAYOUT_NAMES))), None)
BLANK_CANVAS_LAYOUT: int | None = next(
    iter(sorted(_numbers_for({BLANK_CANVAS_NAME}))), None
)

#: Wat er op een layout minstens gevuld moet zijn wil de slide iets zeggen, per NUMMER.
#: Niet alle placeholders: de subtitel (idx 1) is per beleid optioneel, en de tweede kolom
#: van 21/22 hoeft niet gevuld te zijn. Komt uit `LAYOUT_POLICY`, dus op naam.
LAYOUT_REQUIRED: dict[int, frozenset[int]] = {
    number: LAYOUT_POLICY[name].required
    for number, name in LAYOUT_NAMES.items()
    if name in LAYOUT_POLICY and LAYOUT_POLICY[name].required
}


def _xdg_font_dirs() -> list[Path]:
    roots = []
    for entry in (os.environ.get("XDG_DATA_DIRS") or "").split(os.pathsep):
        entry = entry.strip()
        if entry:
            roots.append(Path(entry) / "fonts")
    return roots


def _cloud_font_dirs() -> list[Path]:
    """Office's cloud fonts: `%LOCALAPPDATA%/Microsoft/FontCache/4/CloudFonts/<Familie>/`.

    Dit is waar Montserrat en Lato op een Office-machine werkelijk staan: PowerPoint
    haalt ze als cloud font binnen en zet ze NIET in `C:/Windows/Fonts`. De bestanden
    heten daar `26369665693.ttf` — een nummer, geen familienaam. Naam-matching kan die
    dus per definitie niet vinden, en dat is precies waarom `fonts_measured` leeg bleef
    op een machine waar de render de fonts wél correct zette. De mapnaam draagt de
    familie; de stijl komt uit de nametable (zie `_font_identity`).
    """
    local = os.environ.get("LOCALAPPDATA")
    if not local:
        return []
    return [Path(local) / "Microsoft" / "FontCache" / "4" / "CloudFonts"]


# Every place a brand font can live, Windows and Linux alike. One constant, imported by
# fit_title.py and qa_fit.py, so the two cannot measure with different fonts and then
# disagree about the same title. `assets/fonts` is first: a font shipped with the plugin
# is the same on every machine.
FONT_DIRS: tuple[Path, ...] = tuple(
    [
        SCRIPT_DIR.parent / "assets" / "fonts",
        Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts",
        Path.home() / "AppData/Local/Microsoft/Windows/Fonts",
        Path("/usr/share/fonts"),
        Path("/usr/local/share/fonts"),
        Path.home() / ".local/share/fonts",
        Path.home() / ".fonts",
    ]
    + _cloud_font_dirs()
    + _xdg_font_dirs()
)

FONT_SUFFIXES = ("ttf", "otf", "ttc")


def _normalise_family(name: str) -> str:
    return "".join(c for c in name.lower() if c.isalnum())


def _candidate_files() -> list[Path]:
    """Every font file in FONT_DIRS, once. Cached, because the trees are big."""
    global _CANDIDATE_CACHE
    if _CANDIDATE_CACHE is not None:
        return _CANDIDATE_CACHE
    found: list[Path] = []
    seen: set[str] = set()
    for root in FONT_DIRS:
        if not root.exists():
            continue
        for suffix in FONT_SUFFIXES:
            for pattern in (f"*.{suffix}", f"*.{suffix.upper()}",
                            f"**/*.{suffix}", f"**/*.{suffix.upper()}"):
                for candidate in glob.glob(str(root / pattern), recursive=True):
                    key = candidate.lower()
                    if key in seen:
                        continue
                    seen.add(key)
                    found.append(Path(candidate))
    _CANDIDATE_CACHE = found
    return found


_CANDIDATE_CACHE: list[Path] | None = None
_IDENTITY_CACHE: dict[str, tuple[str, str] | None] = {}


def _font_identity(path: Path) -> tuple[str, str] | None:
    """(familie, stijl) uit de nametable van het fontbestand, of None.

    Nodig voor fontbestanden waarvan de naam niets zegt (Office cloud fonts). Leest via
    Pillow/FreeType; zonder Pillow is er niets te lezen en valt het terug op de naam.
    """
    key = str(path).lower()
    if key in _IDENTITY_CACHE:
        return _IDENTITY_CACHE[key]
    identity: tuple[str, str] | None = None
    try:
        from PIL import ImageFont

        family, style = ImageFont.truetype(str(path), 16).getname()
        identity = (family or "", style or "")
    except Exception:
        identity = None
    _IDENTITY_CACHE[key] = identity
    return identity


def _by_file_name(wanted: str) -> Path | None:
    """De oude, snelle route: de bestandsnaam draagt de familienaam.

    Matching on "contains" alone picks `Montserrat-Black` for `Montserrat`, and Black is
    a good deal wider than Regular — enough to turn a fitting line into a fake overflow.
    So the file whose name adds the least to the family wins, `Regular` counts as
    nothing, and italics only match when the family asks for them.
    """
    best: tuple[int, str] | None = None
    best_path = None
    for candidate in _candidate_files():
        stem = _normalise_family(candidate.stem)
        if wanted not in stem:
            continue
        remainder = stem.replace(wanted, "", 1)
        score = 0 if remainder in ("", "regular") else len(remainder)
        if "italic" in remainder or "oblique" in remainder:
            score += 100
        if best is None or (score, stem) < best:
            best = (score, stem)
            best_path = candidate
    return best_path


def _by_name_table(wanted: str) -> Path | None:
    """De tweede route: familie en stijl uit de nametable van het bestand zelf.

    Alleen gebruikt wanneer de bestandsnaam niets opleverde, en alleen op bestanden
    waarvan het pad de familie noemt (`CloudFonts/Montserrat/…`) — anders zou dit elke
    font op de machine openen.
    """
    best: tuple[int, str] | None = None
    best_path = None
    for candidate in _candidate_files():
        parents = [_normalise_family(p.name) for p in candidate.parents[:2]]
        if not any(part and (part in wanted or wanted in part) for part in parents):
            continue
        identity = _font_identity(candidate)
        if identity is None:
            continue
        family, style = identity
        full = _normalise_family(f"{family} {style}")
        bare = _normalise_family(family)
        if wanted == full or (wanted == bare and _normalise_family(style) in ("", "regular")):
            score = 0
        elif wanted == bare:
            score = 10
        elif wanted in full:
            score = 20 + len(full) - len(wanted)
        else:
            continue
        if "italic" in _normalise_family(style) and "italic" not in wanted:
            score += 100
        if best is None or (score, full) < best:
            best = (score, full)
            best_path = candidate
    return best_path


def find_font_file(family: str) -> Path | None:
    """The font file for a family name like 'Lato Light', or None.

    Twee routes, in deze volgorde: de bestandsnaam (snel, dekt Windows/Linux-installaties
    en `assets/fonts`), en daarna de nametable van bestanden in een map die de familie
    noemt (dekt de Office cloud fonts, waar het bestand een nummer heet). Linux font
    trees are nested, so the search is recursive; `.ttc` collections count too.
    """
    wanted = _normalise_family(family or "")
    if not wanted:
        return None
    return _by_file_name(wanted) or _by_name_table(wanted)


def font_report(families) -> dict:
    """Per familie het gevonden bestand (of None) plus één remediatieregel.

    De QA-scripts melden alleen wélke families gemeten zijn. Dat vertelt niet waarom een
    familie niet gemeten kon worden — geen bestand, of geen Pillow. Dit doet dat wel, en
    het is dezelfde zoekactie, dus de twee kunnen niet uit elkaar lopen.
    """
    try:
        import PIL  # noqa: F401

        pillow = True
    except Exception:
        pillow = False

    found = {name: find_font_file(name) for name in families}
    missing = sorted(name for name, path in found.items() if path is None)
    hint = None
    if not pillow:
        hint = (
            "Pillow ontbreekt — zonder Pillow kan geen enkel font gemeten worden en "
            "vallen alle fit-oordelen terug op een schatting: "
            "pip install -r requirements.txt"
        )
    elif missing:
        hint = (
            f"niet gevonden: {', '.join(missing)} — fit-oordelen op die families zijn "
            "een schatting. Zet de fontbestanden in assets/fonts/ (of installeer ze) "
            "als je echt wilt meten; Gotham is commercieel en staat zelden op een "
            "build-machine."
        )
    return {
        "pillow": pillow,
        "found": {name: (str(path) if path else None) for name, path in found.items()},
        "missing": missing,
        "hint": hint,
    }


@dataclass
class ParagraphStyle:
    """Spacing and indent that decide how tall and how wide a paragraph renders.

    Line spacing and paragraph gaps are given either as a percentage of the line or as
    an absolute number of points; whichever the template used is the one that is set.
    `margin_left_in` is the resolved `marL` — in this template 0.75 in on level 2 and
    1.25 in on level 3, not the 0.25 in per level a naive estimate assumes.
    """

    level: int
    line_spacing_pct: float | None
    line_spacing_pt: float | None
    space_before_pct: float | None
    space_before_pt: float | None
    space_after_pct: float | None
    space_after_pt: float | None
    margin_left_in: float = 0.0
    indent_in: float = 0.0


@dataclass
class RunStyle:
    size_pt: float
    font: str | None
    bold: bool
    italic: bool
    underline: bool
    color: str | None
    highlight: str | None
    size_on_run: bool  # size set on the run itself, not inherited
    color_own: bool  # colour set on the slide, not inherited from layout or master
    font_own: bool  # typeface set on the slide, not inherited from layout or master


def theme_hexes(presentation) -> set[str]:
    """The colours of the themes the slide masters actually use, plus white and black.

    Only the masters' own themes count. A template can carry unused Office themes, and
    those colours are not brand colours.
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    found = set(NEUTRAL_HEX)
    for master in presentation.slide_masters:
        try:
            theme_part = master.part.part_related_by(RT.THEME)
        except KeyError:
            continue
        xml = theme_part.blob.decode("utf-8", "ignore")
        for scheme in re.findall(r"<a:clrScheme.*?</a:clrScheme>", xml, re.S):
            found.update(h.upper() for h in re.findall(r'val="([0-9A-Fa-f]{6})"', scheme))
            found.update(h.upper() for h in re.findall(r'lastClr="([0-9A-Fa-f]{6})"', scheme))
    return found


def _lvl_ppr(lst_style, level: int):
    if lst_style is None:
        return None
    return lst_style.find(f"a:lvl{level + 1}pPr", NS)


def _def_rpr(lst_style, level: int):
    node = _lvl_ppr(lst_style, level)
    if node is None:
        return None
    return node.find("a:defRPr", NS)


def _placeholder_key(shape):
    if not shape.is_placeholder:
        return None
    fmt = shape.placeholder_format
    return (str(fmt.type).split(" ")[0].lower(), fmt.idx)


def _matching_placeholder(container, key):
    if container is None or key is None:
        return None
    wanted_type, wanted_idx = key
    for shape in container.placeholders:
        fmt = shape.placeholder_format
        if fmt.idx == wanted_idx:
            return shape
    if wanted_type in {"title", "center_title"}:
        for shape in container.placeholders:
            if str(shape.placeholder_format.type).split(" ")[0].lower() in {
                "title",
                "center_title",
            }:
                return shape
    return None


def _tx_style_ppr(master, key, level: int):
    if master is None:
        return None
    styles = master._element.find("p:txStyles", NS)
    if styles is None:
        return None
    which = "p:otherStyle"
    if key is not None:
        kind = key[0]
        if kind in {"title", "center_title"}:
            which = "p:titleStyle"
        elif kind in {"body", "subtitle", "object"}:
            which = "p:bodyStyle"
    node = styles.find(which, NS)
    if node is None:
        return None
    return node.find(f"a:lvl{level + 1}pPr", NS)


def _tx_style_rpr(master, key, level: int):
    level_node = _tx_style_ppr(master, key, level)
    return level_node.find("a:defRPr", NS) if level_node is not None else None


def _inheritance(shape):
    """What this shape inherits from: (key, layout placeholder, master placeholder,
    master, presentation defaultTextStyle)."""
    key = _placeholder_key(shape)
    slide = shape.part
    layout_shape = master_shape = master = None
    try:
        slide_layout = slide.slide.slide_layout
        layout_shape = _matching_placeholder(slide_layout, key)
        master = slide_layout.slide_master
        master_shape = _matching_placeholder(master, key)
    except AttributeError:
        pass

    default_style = None
    try:
        default_style = slide.package.presentation_part.presentation._element.find(
            "p:defaultTextStyle", NS
        )
    except AttributeError:
        pass
    return key, layout_shape, master_shape, master, default_style


def _own_lst_style(shape):
    return shape.text_frame._txBody.find("a:lstStyle", NS) if shape.has_text_frame else None


def _chain(shape, paragraph, run, level: int):
    """Every rPr-like element that can define this run's style, nearest first.

    Returns the chain, the placeholder key, and how many leading entries live on the
    slide itself. Anything past that count is inherited from the layout, the master or
    the presentation defaults — template styling, not a choice made on this slide.
    """
    chain = []
    if run is not None and run._r.find("a:rPr", NS) is not None:
        chain.append(run._r.find("a:rPr", NS))

    ppr = paragraph._p.find("a:pPr", NS)
    if ppr is not None:
        end = ppr.find("a:defRPr", NS)
        if end is not None:
            chain.append(end)

    chain.append(_def_rpr(_own_lst_style(shape), level))
    own_count = sum(1 for node in chain if node is not None)

    key, layout_shape, master_shape, master, default_style = _inheritance(shape)

    for candidate in (layout_shape, master_shape):
        if candidate is not None and candidate.has_text_frame:
            chain.append(_def_rpr(candidate.text_frame._txBody.find("a:lstStyle", NS), level))

    chain.append(_tx_style_rpr(master, key, level))

    if default_style is not None:
        chain.append(_def_rpr(default_style, level))

    return [node for node in chain if node is not None], key, own_count


def _first(chain, getter):
    for node in chain:
        value = getter(node)
        if value is not None:
            return value
    return None


def _first_index(chain, getter):
    for index, node in enumerate(chain):
        if getter(node) is not None:
            return index
    return None


def resolve_run(shape, paragraph, run) -> RunStyle:
    level = paragraph.level or 0
    chain, _key, own_count = _chain(shape, paragraph, run, level)

    def latin(node):
        found = node.find("a:latin", NS)
        return found.get("typeface") if found is not None else None

    def srgb(node):
        found = node.find("a:solidFill/a:srgbClr", NS)
        return found.get("val") if found is not None else None

    size = _first(chain, lambda n: n.get("sz"))
    font = _first(chain, latin)
    bold = _first(chain, lambda n: n.get("b"))
    italic = _first(chain, lambda n: n.get("i"))
    underline = _first(chain, lambda n: n.get("u"))
    color = _first(chain, srgb)

    font_index = _first_index(chain, latin)
    color_index = _first_index(chain, srgb)
    highlight = _first(
        chain,
        lambda n: (
            n.find("a:highlight/a:srgbClr", NS).get("val")
            if n.find("a:highlight/a:srgbClr", NS) is not None
            else None
        ),
    )

    run_rpr = run._r.find("a:rPr", NS) if run is not None else None
    return RunStyle(
        size_pt=float(size) / 100 if size else FALLBACK_SIZE_PT,
        font=font,
        bold=bold == "1",
        italic=italic == "1",
        underline=underline not in (None, "none"),
        color=color.upper() if color else None,
        highlight=highlight.upper() if highlight else None,
        size_on_run=run_rpr is not None and run_rpr.get("sz") is not None,
        color_own=color_index is not None and color_index < own_count,
        font_own=font_index is not None and font_index < own_count,
    )


def _ppr_chain(shape, paragraph, level: int):
    """Every pPr-like element that can define this paragraph's spacing, nearest first.

    Same walk as `_chain`, one level up: line spacing and the gap between paragraphs
    live on `a:lvlNpPr`, and in this template they come from the master (`spcPct` 90%
    on the body, `spcBef` of 10pt). Reading them off the paragraph alone gives None.
    """
    chain = [paragraph._p.find("a:pPr", NS), _lvl_ppr(_own_lst_style(shape), level)]

    key, layout_shape, master_shape, master, default_style = _inheritance(shape)
    for candidate in (layout_shape, master_shape):
        if candidate is not None and candidate.has_text_frame:
            chain.append(_lvl_ppr(candidate.text_frame._txBody.find("a:lstStyle", NS), level))
    chain.append(_tx_style_ppr(master, key, level))
    chain.append(_lvl_ppr(default_style, level))
    return [node for node in chain if node is not None]


def _spacing(chain, tag: str) -> tuple[float | None, float | None]:
    """(percentage, points) for `lnSpc`, `spcBef` or `spcAft`; the first one found."""
    for node in chain:
        found = node.find(f"a:{tag}", NS)
        if found is None:
            continue
        pct = found.find("a:spcPct", NS)
        if pct is not None and pct.get("val"):
            return int(pct.get("val")) / 100000, None
        pts = found.find("a:spcPts", NS)
        if pts is not None and pts.get("val"):
            return None, int(pts.get("val")) / 100
    return None, None


def _emu_attr(chain, name: str) -> int | None:
    for node in chain:
        value = node.get(name)
        if value is not None:
            try:
                return int(value)
            except ValueError:
                return None
    return None


def resolve_paragraph(shape, paragraph) -> ParagraphStyle:
    level = paragraph.level or 0
    chain = _ppr_chain(shape, paragraph, level)
    line_pct, line_pt = _spacing(chain, "lnSpc")
    before_pct, before_pt = _spacing(chain, "spcBef")
    after_pct, after_pt = _spacing(chain, "spcAft")
    mar_l = _emu_attr(chain, "marL")
    indent = _emu_attr(chain, "indent")
    return ParagraphStyle(
        level=level,
        line_spacing_pct=line_pct,
        line_spacing_pt=line_pt,
        space_before_pct=before_pct,
        space_before_pt=before_pt,
        space_after_pct=after_pct,
        space_after_pt=after_pt,
        margin_left_in=(mar_l or 0) / EMU_PER_INCH,
        indent_in=(indent or 0) / EMU_PER_INCH,
    )


def _body_pr(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    return shape.text_frame._txBody.find("a:bodyPr", NS)


def resolve_anchor(shape) -> str:
    """The vertical anchor PowerPoint will use, inheritance included.

    A placeholder that `add_slide.py` wrote carries `<a:bodyPr><a:noAutofit/></a:bodyPr>`
    with no anchor, so it takes the layout's. In this template that matters: the column
    heads of layouts 23 and 24 and every idx-1 subtitle are `anchor="ctr"`, and reading
    "t" there puts the text at the top of a box where it actually sits centred.
    """
    own = _body_pr(shape)
    anchor = own.get("anchor") if own is not None else None
    if anchor:
        return anchor

    _key, layout_shape, master_shape, _master, _default = _inheritance(shape)
    for candidate in (layout_shape, master_shape):
        node = _body_pr(candidate) if candidate is not None else None
        if node is not None and node.get("anchor"):
            return node.get("anchor")
    return "t"


def shape_type_name(shape) -> str:
    """`shape_type` als string, ook wanneer python-pptx het type niet kan bepalen.

    Een `<p:sp>` zonder `prstGeom` of `custGeom` laat `shape.shape_type` een
    NotImplementedError gooien ("Shape instance of unrecognized shape type"). PowerPoint
    schrijft die geometrie altijd, maar generatoren niet allemaal — en de scripts lezen
    ook decks die deze plugin niet gebouwd heeft. Eén vreemde vorm mag een QA-run of een
    inspectie niet met een traceback laten eindigen: het antwoord is "onbekend type", en
    de rest van de slide is nog steeds te beoordelen.
    """
    try:
        return str(shape.shape_type or "")
    except NotImplementedError:
        return ""


def role_of(shape) -> str:
    if not shape.is_placeholder:
        return "shape"
    fmt = shape.placeholder_format
    kind = str(fmt.type).split(" ")[0].lower()
    if kind in {"title", "center_title"}:
        return "title"
    if fmt.idx == 1:
        return "subtitle"
    return "body"


def open_deck(path: str | Path):
    return Presentation(str(path))


# A straight apostrophe between two word characters: `risico's`, `'s-Hertogenbosch`
# never gets here because the leading one has no word character in front of it. Only
# the unambiguous case is rewritten; a straight quote used as a quotation mark is left
# alone and reported by qa_typography.py, because turning it into ’ would be wrong.
IN_WORD_APOSTROPHE = re.compile(r"(?<=[^\W_])'(?=[^\W_])", re.UNICODE)
STRAIGHT_QUOTES = re.compile(r"['\"]")


def normalise_apostrophes(text: str) -> str:
    """Rechte apostrof binnen een woord wordt een typografische apostrof (’)."""
    return IN_WORD_APOSTROPHE.sub("’", text)


def deck_positions(unpacked: str | Path) -> dict[str, int]:
    """{'slide3.xml': 1, ...}: the 1-based deck position of every slide file.

    `slideN.xml` is a file name, not a slide number: `add_slide.py --at 1` writes
    slide3.xml at the front of the deck. Every script has to report the position the
    reader sees, or two findings about the same slide carry different numbers.
    """
    from lxml import etree

    root_dir = Path(unpacked)
    presentation = root_dir / "ppt" / "presentation.xml"
    rels = root_dir / "ppt" / "_rels" / "presentation.xml.rels"
    if not presentation.exists() or not rels.exists():
        return {}

    rid_to_slide: dict[str, str] = {}
    for rel in etree.parse(str(rels)).getroot():
        target = rel.get("Target", "")
        if "slideLayout" in (rel.get("Type") or ""):
            continue
        if target.startswith("slides/"):
            rid_to_slide[rel.get("Id")] = target.split("/")[-1]

    xml = presentation.read_text(encoding="utf-8")
    order = re.findall(r'<p:sldId[^>]*r:id="([^"]+)"', xml)
    positions: dict[str, int] = {}
    for index, rid in enumerate(order, start=1):
        name = rid_to_slide.get(rid)
        if name:
            positions[name] = index
    return positions


_UTF8_DONE = False


def use_utf8() -> None:
    """Put stdout AND stderr on UTF-8. Idempotent — reconfigure only ONCE.

    The Windows console defaults to cp1252 and turns every dash, quote and accented
    letter into a question mark, which makes findings hard to read. stderr matters as
    much as stdout here: the scripts raise SystemExit with Dutch messages ("één reeks",
    "— kies uit"), and those go to stderr, where cp1252 mangled them and could even
    raise UnicodeEncodeError on interpolated user text after the XML was already
    written.

    Called on import of this module, so every script in the toolkit emits the same way
    without having to remember it.

    EÉN keer, en dat is de reden voor de vlag. Een tweede `reconfigure()` op een pipe
    op Windows liet de interpreter bij afsluiten struikelen over de oude buffer:
    `Exception ignored in: <_io.TextIOWrapper name='<stdout>'> OSError: [Errno 22]
    Invalid argument`, boven op de JSON. De JSON bleef geldig, maar elke caller die de
    uitvoer parseert kreeg die regels erbij. `emit()` riep dit bij élke aanroep opnieuw
    aan; nu niet meer.
    """
    global _UTF8_DONE
    import sys

    if _UTF8_DONE:
        return
    _UTF8_DONE = True
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8")
        except (AttributeError, OSError, ValueError):
            # pytest replaces these with capture objects that have no reconfigure.
            pass


use_utf8()


def emit(payload) -> None:
    """Print the result as UTF-8 JSON.

    Schrijft in één keer, met een expliciete flush: dan is de JSON weg vóór de
    interpreter afsluit, en kan een halfleeg buffer bij afsluiten geen OSError meer op
    stdout achterlaten naast de JSON.
    """
    import json
    import sys

    text = json.dumps(payload, ensure_ascii=False, indent=2) + "\n"
    try:
        sys.stdout.write(text)
        sys.stdout.flush()
    except (OSError, ValueError):
        # Pipe al dicht: er is niets meer te melden en een traceback hierover helpt
        # niemand.
        pass


def path_is_writable(target: str | Path, *, via_tempfile: bool = False) -> bool:
    """Kan er straks op `target` geschreven worden? Te toetsen VÓÓR het werk, niet erna.

    Dit is het beleid dat de README onder *Referentie* opschrijft: een geïnstalleerde
    plugin staat in een read-only map, en een script dat dat pas op de laatste
    schrijfactie merkt heeft dan al een sjabloon uitgepakt en omgezet — al het werk
    gedaan voor een kale `PermissionError` die niet zegt wat je moet doen. Elk script dat in de
    pluginmap schrijft toetst zijn doel dus vooraf en handelt er zelf naar:
    `layout_catalog.py` wijkt uit (cwd, dan stdout), `prune_template.py` stopt met een
    melding die naar `--out` verwijst omdat een sjabloon dat stilletjes ergens anders
    landt niemand helpt.

    Openen in plaats van `os.access`: `os.access` kijkt op Windows alleen naar het
    read-only-attribuut en negeert de ACL's, en juist een DENY-ACL is waarmee een
    installatie dichtgezet wordt. Alleen een echte schrijfpoging is hier betrouwbaar.

    `via_tempfile=True` als het schrijven via `helpers.rezip()` gaat. Die schrijft naar
    een tempfile NAAST het doel en `os.replace`t die eroverheen, dus daar moet de MAP
    beschrijfbaar zijn — een toets op alleen het bestand zou precies die faalwijze
    missen. Zonder de vlag volstaat het bestand zelf (`write_text` op een bestaand
    bestand in een verder dichte map werkt wél), en dat is wat `layout_catalog.py` doet.
    """
    target = Path(target)
    try:
        target.parent.mkdir(parents=True, exist_ok=True)
        if target.exists() and not via_tempfile:
            with target.open("a", encoding="utf-8"):
                return True
        probe = target.parent / f".{target.name}.probe"
        probe.write_text("", encoding="utf-8")
        probe.unlink()
        if target.exists():
            with target.open("a", encoding="utf-8"):
                pass
        return True
    except (OSError, ValueError):
        # ValueError: een pad met een NUL erin komt niet eens bij het bestandssysteem.
        # Dat is geen schrijfbaar doel, en een traceback daarover helpt niemand.
        return False


def iter_text(slide):
    """Yield (shape, paragraph, run) for every run with text on the slide."""
    for shape in slide.shapes:
        for item in _iter_shape_text(shape):
            yield item


def _iter_shape_text(shape):
    if shape_type_name(shape).startswith("GROUP"):
        for child in shape.shapes:
            yield from _iter_shape_text(child)
        return
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        yield shape, paragraph, run
        return
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            yield shape, paragraph, run
