"""Tellingen op een gebouwd deck. Zes drempels, vier cijfers, geen vormbeoordeling.

Usage:
    python qa_tellingen.py <deck.pptx>
    python qa_tellingen.py <deck.pptx> --renders <map>   # ook de registerverdeling
    python qa_tellingen.py <deck.pptx> --strict          # ook een warn geeft exit 1

Dit is een hygiënerapport in dezelfde categorie als `qa_text.py`, en **geen tweede
poort.** De poort van deze plugin is de outline; daarna beslist de render, met
`deck-visual-reviewer` als het oog. Dat staat zo in de README ("Poorten: één: de
outline") en dit script verandert dat niet. De verantwoording hoort hier te staan omdat
dit precies de plek is waar deze plugin anders de vormgevingspolitie herbouwt die hij
bewust weglaat: een script dat compositie telt, groeit uit zichzelf naar een script dat
compositie afkeurt, en dan vermijdt de bouwer regels in plaats van slides te maken.
Daarom staat de grens hier hard:

* **`critical` blokkeert**, en alleen wat mechanisch en zonder interpretatie vast te
  stellen is én afwijkt van een besluit dat de bouwer zelf heeft genomen. Dat zijn er
  drie: meer dan één maat per rol (de vier maten uit de outline, `vormentaal.md` §2),
  twee letterfamilies in dezelfde alinea (§9) en de hoge punt binnen een regel (§9).
* **`warn` is een aanwijzing**, geen afkeuring: bandfrequentie, nul exhibits in een deck
  met cijfers, en de maatsprong per slide. Alle drie zijn ze te meten, maar geen van de
  drie is zonder de render te beoordelen — een prozaslide met een gekleurde kop haalt
  een lage maatsprong en kan de sterkste slide van de deck zijn.
* **Cijfers zonder oordeel**: woorden per slide en per element, de registerverdeling en
  de herhaalde plattegrond. Die krijgen geen severity, ook niet als ze slecht uitvallen.

Waarom die laatste categorie bestaat, en waarom dat zo moet blijven: een `critical` op
woorden per slide leert de bouwer tekst versnipperen in plaats van reduceren, en dat is
precies het defect waar deze werklijn vandaan komt (het gemeten deck: gemiddeld 177
woorden per contentslide, piek 255, vier kolommen met elk zeventig woorden). Wie deze
telling later tot drempel promoveert, krijgt tien slides van negentig woorden terug in
plaats van vijf slides met een schema. Hetzelfde geldt voor de registerverdeling en de
plattegrond: die vragen om een besluit over de deck als geheel, en dat besluit hoort in
de outline en op de render, niet in een exitcode.

Wat hier NIET in zit en er ook niet in komt: `qa_fit.py` en `qa_typography.py`. Die
werden door vijf scripts aangehaald als "de poort in QA-only-modus" en hebben nooit in
deze repo bestaan; ze zijn opgeruimd in plaats van geschreven.

De drempels zijn nagemeten, niet gekozen. Alle getallen hieronder komen uit de twee
fixtures van de nulmeting (23 en 19 slides, `Werksessie 1 impactmeten`) en uit de
getallen die al in `vormentaal.md` staan:

| telling | niveau | drempel | gemeten op de twee fixtures |
|---|---|---|---|
| maten per rol | deck | één maat per rol | de body op 14pt in beide decks (§2 zegt 16); de Montserrat SemiBold-rollen op vijf maten in het spreekdeck — label 13 en 14pt, kop 15, 16 en 18pt |
| bandfrequentie | deck | ten hoogste één per vier slides (§10) | 14 van de 17 en 12 van de 13 contentslides, tegen een ruimte van 4 en 3 |
| exhibits bij cijfers | deck | minstens één grafiek, tabel, schema of verdeling (§12) | nul grafieken en nul tabellen bij 162 en 85 cijferfeiten |
| maatsprong | slide | grootste eigen maat / kleinste >= 2 (§1) | 1,00 tot 2,91 (13 van 17 slides eronder) en 1,15 tot 3,64 (9 van 13); de afgekeurde deck uit §1 haalde 1,36, de referentie 3 tot 5 |
| twee families in één alinea | slide | nul (§9) | 71 en 40 alinea's |
| hoge punt binnen een regel | slide | nul (§9) | 8 en 3 regels |

Op een schone testbouw (cover, twee gecomponeerde contentslides met een tabel, outro,
gebouwd met `shapes.py` en `add_table.py`) staat dit script op nul `critical` en één
`warn`: de maatsprong op de slide waar een navy paneel en een afsluitband de hiërarchie
dragen in plaats van de maat. Dat is precies de reden dat de maatsprong een `warn` is en
geen `critical`.

Output is compacte JSON. Exit 1 zodra er een `critical` staat, of bij `--strict` ook op
een `warn`.
"""

from __future__ import annotations

import argparse
import re
import sys
from collections import Counter, defaultdict
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

from _deck import emit  # noqa: E402
from shapes import DRAGER_VLOER, HUE  # noqa: E402

# Eén definitie voor "wat is chrome", "wat is een titel" en "op welke layouts". Deze
# komen uit `qa_text.py` in plaats van hier opnieuw te staan: twee QA-scripts die het
# oneens zijn over of een subregel chrome is, is een bug die je niet ziet.
from qa_text import (  # noqa: E402
    HEADLINE_BY_IDX_LAYOUTS,
    INTENTIONALLY_BLANK_LAYOUTS,
    is_chrome,
    is_title,
    layout_number,
    shape_key,
    walk,
)

EMU_PER_INCH = 914400

#: De hoge punt en zijn lookalikes, tussen twee stukken tekst op dezelfde regel. Dit
#: kijkt op de slide-XML en nooit op documentatie: `reference/sjabloon.md` gebruikt de
#: hoge punt legitiem als coördinaatscheiding in een tabel ("3,2 · 1,4"), en een grep
#: over de repo zou daar vallen.
HOGE_PUNT = re.compile(r"\S[ \t]*[·•∙⋅‧][ \t]*\S")

#: Een cijferfeit: een getal met of zonder eenheid. Jaartallen en paginanummers vallen
#: erbuiten (chrome wordt niet gelezen, en vier cijfers achter elkaar is geen bedrag
#: zonder scheidingsteken). Dit bepaalt alleen of een deck "cijfers draagt".
CIJFER = re.compile(
    r"(?<![\w/])(?:€\s*)?\d{1,3}(?:[.\s]\d{3})*(?:,\d+)?\s*"
    r"(?:%|procent|mln|mld|fte|uur|dagen|weken|maanden|jaar)?"
)

#: Grens tussen bodymaat en voetnootmaat, en tussen body en de dichte variant. §2 zet
#: 16pt Lato Light als body, 12pt als de dichte variant voor een kaartenrij van drie of
#: meer of een tabelcel, en 11pt als voetnoot. Die drie zijn dus eigen rollen en geen
#: drie maten voor dezelfde rol — anders meldt dit script een kaartenrij als defect.
VOETNOOT_PLAFOND = 11.5
DICHT_PLAFOND = 12.5

#: Hetzelfde plafond, maar voor de Montserrat SemiBold-kant. Dit stond er niet en dat was
#: een fout: de drie maatbanden hierboven golden alleen voor Lato, en élke Montserrat
#: SemiBold tot en met 14pt viel in één rol `label`. Nagemeten op `werk-b/build/maatstaf`,
#: de vier voorbeeldslides uit `assets/maatstaf/11`-`14`: daar staat een kapitaallabel van
#: 11pt boven een tabelkolom (`Kolomkop 1..3` op `12`, `Askop links` en `As 1..4` op `13`),
#: een rijkop en een kaartlabel van 12pt (`Rij 1 rol` op `12`, `Kaart 1 label` op `11`) en
#: een kapitaallabel van 14pt (`Kolom 1 kop` op `14`). Met één rol `label` meldde dit
#: script daarop `critical: de rol 'label' staat op 3 maten` — op de lat zelf, en de bouwer
#: heeft daarop besluit 2 teruggebracht naar één labelmaat van 12pt. Dat besluit nam het
#: telscript en niet de bouwer, en dat is precies wat dit script niet mag doen.
#:
#: §2 geeft de drie maten elk hun eigen werk: "12pt is de dichte variant voor een
#: kaartenrij van drie of meer of een tabelcel, en de vloer. 14pt is het kapitaallabel", en
#: §11 schrijft de bronregel op 11pt voor. Drie maten, drie rollen. Wat overblijft is de
#: toets die wél moet vuren: dezelfde rol op twee maten bínnen dezelfde band — in het
#: afgekeurde deck staat het kapitaallabel op 13 én 14pt, en dat is één rol op twee maten.
LABEL_PLAFOND = 14.5

#: De bodymaat die §2 voorschrijft. Afwijken mag, maar dan is het een besluit: in beide
#: gemeten decks zakte de body ongemerkt naar 14pt zonder dat iemand dat koos.
BODY_NORM = 16.0

#: Een afsluitband: een gevuld vlak over vrijwel de volle zonebreedte, laag, en tegen de
#: onderkant. §10 meet hem als 12,52 bij 1,25 in met één regel erin. De hoogtevloer van
#: 0,30 in houdt lijnwerk en haarlijnen erbuiten: die zijn geen band.
BAND_BREEDTE_MIN = 10.0
BAND_HOOGTE_MAX = 1.75
BAND_ONDERKANT_MIN = 5.60
BAND_GEVULD_HOOGTE_MIN = 0.30

#: Een badge: een kleine ronde vorm met alleen een kort getal erin. `punt()` in `shapes.py`
#: schrijft precies dat — `prst="ellipse"` (of een `pie` over een lege punt), vierkant,
#: `insets=(0,0,0,0)`, `anchor="ctr"`, en als tekst één gecentreerde run. Dat geeft een
#: aangrijpingspunt dat niet te misbruiken is: de vorm moet rond zijn, vierkant, klein, en
#: er mag niets anders in staan dan één tot drie cijfers. "Alles wat kort is" zou hier de
#: verkeerde regel zijn — een kolomkop van twee letters is nog steeds een kop.
#:
#: Gemeten: `assets/maatstaf/11` zet de badge op 0,50 in, `leavebehind` op 0,44 in, en
#: `maatstaf/12` gebruikt dezelfde vorm van 0,16 in als puntenmeter zonder tekst.
BADGE_ZIJDE_MAX = 0.80
BADGE_VERHOUDING = (0.75, 1.33)
BADGE_PRESETS = {"ellipse", "pie", "circle"}
BADGE_TEKST = re.compile(r"^\d{1,3}[.)]?$")

#: De rollen die per definitie een BAND hebben en dus geen enkele maat: §2 zet de drager
#: op 28 tot 40pt, en Montserrat Light onder die vloer is een citaat. Een deck dat één
#: drager op 32 en de afsluitende op 40pt zet, is niet inconsistent maar precies wat de
#: band toestaat; `qa_text.py` waakt daar al over de vloer, het plafond en de frequentie.
#: Hun maten staan wél in de JSON, alleen zonder `critical`.
BAND_ROLLEN = {"drager", "citaat"}

#: Maatsprong: grootste eigen maat gedeeld door kleinste, per slide (§1).
MAATSPRONG_VLOER = 2.0

#: Aantal cijferfeiten waarboven een deck "cijfers draagt". Gemeten: de twee fixtures
#: staan op 166 en 87, en een deck met een datum op de cover en een paginanummer haalt
#: er nul tot drie. Zes is de grens waaronder een enkel getal in een zin geen exhibit
#: verdient.
CIJFERS_VLOER = 6

#: De accentslots. Een accent in de letter is de andere manier om hiërarchie te zetten
#: (§1), en dat feit staat per slide in de JSON zodat een lage maatsprong te wegen is.
ACCENTSLOTS = {HUE[naam] for naam in ("oranje", "grapefruit", "royal", "sky", "emerald")}

#: Registerklassen op de render, zoals de nulmeting ze meet. Ruwe methode, en niet
#: één-op-één vergelijkbaar met de percentages in §5 — de waarde zit in de vergelijking
#: vóór en ná, met hetzelfde script.
WIT_SAT_MAX = 0.06
WIT_VAL_MIN = 0.93
VERZADIGD_SAT_MIN = 0.25
VERZADIGD_VAL_MAX = 0.55


def add(findings, slide, check, message, severity="warn"):
    findings.append(
        {"slide": slide, "check": check, "severity": severity, "message": message}
    )


def inches(shape) -> tuple[float, float, float, float] | None:
    """(x, y, w, h) in inch, of None als de vorm geen geometrie heeft."""
    try:
        return (
            shape.left / EMU_PER_INCH,
            shape.top / EMU_PER_INCH,
            shape.width / EMU_PER_INCH,
            shape.height / EMU_PER_INCH,
        )
    except (AttributeError, TypeError):
        return None


def is_contentslide(layout: int | None) -> bool:
    return layout not in HEADLINE_BY_IDX_LAYOUTS | INTENTIONALLY_BLANK_LAYOUTS


def schemeclr(run) -> str | None:
    """Het themaslot van een tekstrun (`dk2`, `accent1`, ...), of None als hij erft."""
    rpr = run._r.find(qn("a:rPr"))
    if rpr is None:
        return None
    fill = rpr.find(qn("a:solidFill"))
    if fill is None:
        return None
    scheme = fill.find(qn("a:schemeClr"))
    return scheme.get("val") if scheme is not None else "hex"


def vol_vlak_met_tekst(shape) -> bool:
    """Draagt deze vorm een volle vulling met tekst erin?

    Dit is de tweede helft van de OF-toets uit §1. Hiërarchie mag uit maat komen, of uit
    gewicht en kleur — en "kleur" is niet alleen een gekleurde letter. Een verzadigd
    rijlabel met wit erin doet hetzelfde werk, en dat is precies wat de sterkste
    voorbeeldslide van de veertien doet (`assets/maatstaf/12`: vier volle rijlabels tegen
    een nauwelijks getint paneel). Zonder deze uitweg meldt dit script die compositie als
    "geen hiërarchie", en dat is het omgekeerde van wat er staat.

    Een lichte container telt niet mee: die is achtergrond, geen drager. §4 zet containers
    op alpha 6000 tot 14000, dus alles onder de helft is tint en geen volle vulling.
    """
    try:
        if shape.fill.type != 1:
            return False
    except (AttributeError, TypeError, ValueError):
        return False
    if not (getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()):
        return False
    sppr = shape._element.find(qn("p:spPr"))
    fill = sppr.find(qn("a:solidFill")) if sppr is not None else None
    if fill is None:
        return False
    for kleur in fill:
        alpha = kleur.find(qn("a:alpha"))
        if alpha is not None:
            try:
                return int(alpha.get("val", "100000")) >= 50000
            except (TypeError, ValueError):
                return True
    return True


def is_badge(shape) -> bool:
    """Is deze vorm een genummerde badge, en dus een merkteken in plaats van een tekstrol?

    `punt()` in `shapes.py` schrijft precies één vorm: rond (`prst="ellipse"`, of een
    `pie` over een lege punt voor de halve variant), vierkant, klein, insets nul,
    `anchor="ctr"`, en als inhoud één gecentreerde run. Die vijf kenmerken samen zijn een
    aangrijpingspunt dat niet te misbruiken is. "Alles wat kort is" zou de verkeerde regel
    zijn: een kolomkop van twee letters is nog steeds een kop.

    Waarom dit nodig is: `rol_van()` zag Montserrat SemiBold op 16pt in een badge en telde
    hem als tweede koprol naast de koppen van 18pt. De bouwer van een testdeck heeft
    daarop álle badges naar 18pt gezet om de telling te halen — op een badge van 0,44 in
    is dat aan de grote kant. Een badge codeert een rangnummer en leest als vorm, niet als
    tekst; hij hoort dus niet in de rollentelling.

    Gemeten: `assets/maatstaf/11` zet de badge op 0,50 in, de testdeck op 0,44 in.
    """
    geo = inches(shape)
    if not geo:
        return False
    _, _, w, h = geo
    if max(w, h) > BADGE_ZIJDE_MAX or not h:
        return False
    if not BADGE_VERHOUDING[0] <= w / h <= BADGE_VERHOUDING[1]:
        return False
    prst = None
    sppr = shape._element.find(qn("p:spPr"))
    geom = sppr.find(qn("a:prstGeom")) if sppr is not None else None
    if geom is not None:
        prst = geom.get("prst")
    if prst not in BADGE_PRESETS:
        return False
    if not getattr(shape, "has_text_frame", False):
        return False
    tekst = shape.text_frame.text.strip()
    return bool(tekst) and bool(BADGE_TEKST.match(tekst))


def rol_van(font: str | None, pt: float | None, in_tabel: bool) -> str | None:
    """De rol waar deze run in valt, of None als er niets over te zeggen is.

    De rol volgt uit familie plus maat, want dat is wat §2 vastlegt: drager in
    Montserrat Light van 28 tot 40pt, kop in Montserrat SemiBold op 18pt, body in Lato
    Light op 16pt, voetnoot op 11pt. Erft de run zijn font of zijn maat, dan staat er
    niets in de XML en wordt er niets geteld — een geërfde placeholderrun is geen
    besluit van de bouwer.

    Tabelcellen en de dichte 12pt-variant krijgen hun eigen rol. Zonder die splitsing
    meldt dit script een kaartenrij van drie op 12pt naast een prozakolom op 16pt als
    "twee bodymaten", en dat is precies wat §2 toestaat.

    **En de Montserrat-kant krijgt dezelfde drie banden als de Lato-kant, want dat stond
    er niet en het heeft een ontwerpbesluit gekost.** Élke Montserrat SemiBold tot en met
    14pt viel in één rol `label`, terwijl §2 en §11 die maten elk hun eigen werk geven:
    11pt is de bronregel, 12pt is de dichte variant voor een kaartenrij van drie of meer
    en voor een tabelcel, 14pt is het kapitaallabel. Op de lat zelf — de vier
    voorbeeldslides uit `assets/maatstaf/11`-`14` — staan alle drie naast elkaar, en dit
    script meldde daarop `critical: de rol 'label' staat op 3 maten`. De bouwer van de
    testdeck heeft daarop besluit 2 teruggebracht naar één labelmaat van 12pt. Dat besluit
    nam het telscript en niet de bouwer.

    Wat overblijft is de toets die wél moet vuren: dezelfde rol op twee maten bínnen
    dezelfde band. In het afgekeurde deck staat het kapitaallabel op 13 én 14pt, en dat
    is één rol op twee maten.
    """
    if not font or pt is None:
        return None
    familie = font.strip().lower()
    if in_tabel:
        return "tabelcel"
    if familie.startswith("montserrat"):
        if "light" in familie:
            return "drager" if pt >= DRAGER_VLOER else "citaat"
        if pt > LABEL_PLAFOND:
            return "kop"
        if pt <= VOETNOOT_PLAFOND:
            return "labelvoetnoot"
        if pt <= DICHT_PLAFOND:
            return "labeldicht"
        return "label"
    if familie.startswith("lato"):
        if pt <= VOETNOOT_PLAFOND:
            return "voetnoot"
        if pt <= DICHT_PLAFOND:
            return "dicht"
        return "body"
    return None


def band_shapes(slide) -> list[str]:
    """Namen van de afsluitbanden op deze slide: alleen de gevulde vorm.

    Een band is een **gevuld** vlak over vrijwel de volle zonebreedte tegen de onderkant.
    §10 meet hem als 12,52 bij 1,25 in met één regel erin, en het is de vulling die hem
    luid maakt: "een band van 12,52 bij 1,25 in met één regel erin is dezelfde vorm, of
    hij navy, oranje-tint of warmgrijs is".

    **Hier stond eerst een tweede vorm bij, en dat was een fout die het ontwerp heeft
    aangestuurd.** De eerdere versie telde ook "een tekstblok met een volle-breedte lijn
    er direct boven", op de gedachte dat dat hetzelfde leest. Maar §10 biedt juist dát
    aan als uitweg uit de band: "de sluitregel op wit: één regel zonder vulling, aanhef
    in Lato Semibold, hooguit een streep erboven". Het aanbevolen alternatief was dus
    zelf een band volgens de teller. Nagemeten gevolg op een deck dat net met deze skill
    was gebouwd: 12 van de 17 contentslides "sloten af met een band" terwijl er drie
    gevulde banden op stonden, en de bouwer heeft daarop een merkstreepje van 12,52 naar
    1,60 in ingekort om de teller tevreden te stellen. Dat is een vormbesluit dat door een
    telling is ingegeven en niet door de render, en dat is precies wat dit script niet mag
    doen.

    Wat deze telling daardoor niet meer ziet: een deck dat de band vervangt door tien keer
    dezelfde sluitregel onder dezelfde streep. Dat is eenvormigheid, en die hoort ook thuis
    waar de rest van de eenvormigheid ligt — bij de plattegrond in de outline en bij het
    oog van `deck-visual-reviewer`, die er expliciet op let. Een teller die het verschil
    tussen een afsluiter en de aanbevolen uitweg niet kan zien, hoort het niet te tellen.
    """
    vormen = [(shape, inches(shape)) for shape in walk(slide.shapes)]
    vormen = [(shape, geo) for shape, geo in vormen if geo]

    banden = []
    for shape, (x, y, w, h) in vormen:
        if is_chrome(shape):
            continue
        if not (w >= BAND_BREEDTE_MIN and BAND_GEVULD_HOOGTE_MIN <= h <= BAND_HOOGTE_MAX):
            continue
        if y + h < BAND_ONDERKANT_MIN:
            continue
        try:
            gevuld = shape.fill.type == 1
        except (AttributeError, TypeError, ValueError):
            gevuld = False
        if gevuld:
            banden.append(getattr(shape, "name", "?"))
    return banden


def plattegrond(slide, layout: int | None, banden: int) -> str:
    """De plattegrond van de slide in één sleutel, mechanisch benaderd.

    §10 vraagt de plattegrond in vier woorden ("drie kaarten, open onderkant") en zegt:
    komt één plattegrond meer dan twee keer voor, dan is er een slide die opnieuw
    ontworpen moet worden. Dat oordeel hoort in de outline. Wat een script wél kan, is
    de geometrie samenvatten: hoeveel eigen tekstblokken, in hoeveel kolommen en rijen,
    met of zonder band en exhibit. Twee slides met dezelfde sleutel hebben dezelfde
    plattegrond; twee slides met een verschillende sleutel kunnen er nog steeds
    hetzelfde uitzien, dus dit is een ondergrens en geen bewijs.
    """
    kolommen: set[int] = set()
    rijen: set[int] = set()
    blokken = 0
    exhibits = 0
    for shape in walk(slide.shapes):
        if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
            exhibits += 1
            continue
        if not getattr(shape, "has_text_frame", False) or is_chrome(shape):
            continue
        if not shape.text_frame.text.strip() or is_title(shape, layout):
            continue
        geo = inches(shape)
        if geo is None:
            continue
        x, y, w, h = geo
        if w >= BAND_BREEDTE_MIN and h <= BAND_HOOGTE_MAX:
            continue  # de band staat al in de sleutel
        blokken += 1
        kolommen.add(round(x * 2))  # halve inch: fijner en het is ruis
        rijen.add(round(y * 2))
    return (f"{blokken} blokken / {len(kolommen)} kolommen / {len(rijen)} rijen"
            f" / {banden} band / {exhibits} exhibit")


def registers(map_pad: Path) -> dict:
    """Aandeel wit, tint en verzadigd per gerenderde slide.

    Alleen uit een render te halen, niet uit de XML: het aandeel hangt af van hoe groot
    een vlak gezet wordt en hoeveel wit ernaast blijft staan, en dat weet je pas als het
    getekend is. Daarom `--renders` en niet standaard. Zonder Pillow zegt dit veld dat
    er niet gemeten is; het doet niet alsof de verdeling in orde is.

    Zelfde klassering als de nulmeting, zodat de vergelijking vóór en ná op dezelfde
    methode staat: wit is verzadiging < 6 procent en waarde > 93 procent, verzadigd is
    verzadiging > 25 procent of waarde < 55 procent, de rest is tint. De PNG wordt eerst
    naar 240 px breed gebracht — de verhouding verandert daar niet meetbaar van en het
    scheelt een factor honderd rekenwerk.
    """
    try:
        from PIL import Image
    except ImportError:
        return {"gemeten": False,
                "hint": "Pillow ontbreekt, dus de registerverdeling is niet gemeten. "
                        "Dat zegt niets over de verdeling zelf — er is niet gekeken. "
                        "`pip install pillow` en opnieuw met --renders."}

    pngs = sorted(map_pad.glob("*.png"))
    if not pngs:
        return {"gemeten": False,
                "hint": f"geen PNG's in {map_pad} — render eerst met render.py"}

    wit_grens = WIT_SAT_MAX * 255, WIT_VAL_MIN * 255
    vol_grens = VERZADIGD_SAT_MIN * 255, VERZADIGD_VAL_MAX * 255

    per_slide = []
    for png in pngs:
        with Image.open(png) as beeld:
            klein = beeld.convert("RGB")
            klein = klein.resize((240, max(1, round(240 * klein.height / klein.width))))
            ruw = klein.convert("HSV").tobytes()
        sats, vals = ruw[1::3], ruw[2::3]
        totaal = len(sats)
        wit = verzadigd = 0
        for sat, val in zip(sats, vals):
            if sat < wit_grens[0] and val > wit_grens[1]:
                wit += 1
            elif sat > vol_grens[0] or val < vol_grens[1]:
                verzadigd += 1
        per_slide.append({
            "render": png.name,
            "wit": round(100 * wit / totaal),
            "tint": round(100 * (totaal - wit - verzadigd) / totaal),
            "verzadigd": round(100 * verzadigd / totaal),
        })
    witten = sorted(r["wit"] for r in per_slide)
    vollen = sorted(r["verzadigd"] for r in per_slide)
    return {
        "gemeten": True,
        "methode": "wit s<6% v>93% / verzadigd s>25% of v<55%, op 240 px breed",
        "bandbreedte": {"wit": [witten[0], witten[-1]],
                        "verzadigd": [vollen[0], vollen[-1]]},
        "ijkpunt": "§5 meet de referentie op bijna helemaal wit (85 tot 88 procent) OF "
                   "echt verzadigd (20 tot 37 procent), en geen enkele slide in het "
                   "midden; de afgekeurde deck lag met 44 tot 53 procent wit over de hele "
                   "lengte in dezelfde middenband. De geërfde fotodividers (80 procent "
                   "verzadigd) en de oranje outro (99 procent) zeggen niets over de "
                   "compositie — lees deze reeks op de contentslides.",
        "per_slide": per_slide,
    }


def analyse(deck: Path, renders: Path | None = None) -> dict:
    presentatie = Presentation(str(deck))
    findings: list[dict] = []

    maten_per_rol: dict[str, dict[float, list[int]]] = defaultdict(
        lambda: defaultdict(list))
    woorden_per_slide: dict[int, int] = {}
    woorden_per_element: list[int] = []
    bandslides: list[int] = []
    plattegronden: Counter = Counter()
    plattegrond_slides: dict[str, list[int]] = defaultdict(list)
    sprongen: dict[int, dict] = {}
    charts = tables = cijfers = contentslides = 0

    for nummer, slide in enumerate(presentatie.slides, start=1):
        layout = layout_number(slide)
        content = is_contentslide(layout)
        contentslides += 1 if content else 0
        titels = {shape_key(s) for s in walk(slide.shapes) if is_title(s, layout)}

        eigen_maten: list[float] = []
        accent_in_letter = False
        vol_vlak = False
        woorden = 0

        for shape in walk(slide.shapes):
            if getattr(shape, "has_chart", False):
                charts += 1
            if getattr(shape, "has_table", False):
                tables += 1
            if not is_chrome(shape) and vol_vlak_met_tekst(shape):
                vol_vlak = True

        for shape in walk(slide.shapes):
            in_tabel = False
            frames = []
            if getattr(shape, "has_table", False):
                in_tabel = True
                frames = [cel.text_frame for rij in shape.table.rows for cel in rij.cells]
            elif getattr(shape, "has_text_frame", False) and not is_chrome(shape):
                frames = [shape.text_frame]
            if not frames:
                continue

            in_titel = shape_key(shape) in titels
            badge = is_badge(shape)
            element_woorden = 0

            for frame in frames:
                for para in frame.paragraphs:
                    families: set[str] = set()
                    regel = ""
                    for run in para.runs:
                        tekst = run.text
                        regel += tekst
                        element_woorden += len(tekst.split())
                        font = run.font.name
                        pt = run.font.size.pt if run.font.size is not None else None
                        if font:
                            families.add(font.strip().split()[0].lower())
                        if not in_titel:
                            if schemeclr(run) in ACCENTSLOTS:
                                accent_in_letter = True
                            cijfers += len(CIJFER.findall(tekst))
                            rol = rol_van(font, pt, in_tabel)
                            if rol and tekst.strip() and not badge:
                                maten_per_rol[rol][pt].append(nummer)
                            if pt is not None and tekst.strip():
                                eigen_maten.append(pt)

                    if {"montserrat", "lato"} <= families:
                        add(findings, nummer, "twee-families-in-alinea",
                            "Montserrat en Lato staan in dezelfde alinea: "
                            f'"{regel.strip()[:60]}". Eén regel is één familie '
                            "(vormentaal §9) — de aanhef binnen een doorlopende regel is "
                            "Lato Semibold, en Montserrat SemiBold staat op wat lósstaat "
                            "en op zijn eigen regel begint. `para()` in shapes.py weigert "
                            "deze mix bij het bouwen; komt hij hier langs, dan is deze "
                            "slide op een andere manier gemaakt.", "critical")
                    if not in_titel and HOGE_PUNT.search(regel):
                        add(findings, nummer, "hoge-punt-in-regel",
                            "de hoge punt scheidt twee feiten op één regel: "
                            f'"{regel.strip()[:60]}". Twee feiten zijn twee regels, twee '
                            "cellen of twee elementen (vormentaal §9), ook in een label "
                            "en in een bronregel.", "critical")

            woorden += element_woorden
            if element_woorden:
                woorden_per_element.append(element_woorden)

        banden = band_shapes(slide)
        if banden:
            bandslides.append(nummer)

        if content:
            woorden_per_slide[nummer] = woorden
            sleutel = plattegrond(slide, layout, len(banden))
            plattegronden[sleutel] += 1
            plattegrond_slides[sleutel].append(nummer)

            if len(eigen_maten) >= 1:
                sprong = max(eigen_maten) / min(eigen_maten)
                sprongen[nummer] = {
                    "sprong": round(sprong, 2),
                    "klein": min(eigen_maten),
                    "groot": max(eigen_maten),
                    "accent_in_letter": accent_in_letter,
                    "vol_vlak_met_tekst": vol_vlak,
                }
                # §1 stelt één toets met een OF erin: is de drager minstens tweeënhalf
                # keer de bodymaat, óf onderscheidt hij zich in gewicht en kleur? Deze
                # check las daar een EN van, en dat maakte hem onhaalbaar met de vier
                # maten die §2 zelf voorschrijft. Nagemeten op twee decks die net met de
                # skill waren gebouwd: 23 warns, en dat waren álle warns die er waren.
                # Zonder drager van 28pt of meer is de grootste eigen maat de kop van
                # 18pt en de kleinste de voetnoot van 11pt, dus de sprong is 1,64 en
                # lager — en tegelijk mag ten hoogste één slide op drie een drager van
                # 28pt of groter dragen (§1). De check eiste dus op twee derde van de
                # slides precies wat §1 daar verbiedt.
                #
                # Wat overblijft is het defect dat hij bedoelde te vinden: een slide
                # zonder maatsprong én zonder kleur in de letter. Dat is de afgekeurde
                # deck, letterlijk — maatsprong 1,36, "31 gevulde vlakken en nul
                # gekleurde letters" (§3). Staat er wel een accent in de letter, dan
                # draagt dat de hiërarchie en is er niets te melden.
                # `vol_vlak` staat wél in de JSON maar is géén uitweg, en dat is
                # nagemeten: het afgekeurde deck uit §3 had "31 gevulde vlakken en nul
                # gekleurde letters" en maatsprong 1,36. Een volle vulling als excuus
                # laat juist dát deck ongemoeid — geprobeerd, en toen viel de melding op
                # het afgekeurde deck naar nul terwijl twee net gebouwde decks er negen
                # kregen. Precies verkeerd om. Een volle vulling die de drager draagt en
                # een volle vulling als behang zijn mechanisch niet te scheiden, dus
                # meldt dit script het en weegt de bouwer het op de render.
                if sprong < MAATSPRONG_VLOER and not accent_in_letter:
                    add(findings, nummer, "maatsprong",
                        f"maatsprong {sprong:.2f}: de grootste eigen maat is "
                        f"{max(eigen_maten):.0f}pt en de kleinste {min(eigen_maten):.0f}pt "
                        f"(vloer {MAATSPRONG_VLOER:.0f}, vormentaal §1; de afgekeurde deck "
                        "haalde 1,36, de referentie 3 tot 5), én er staat géén accent in "
                        "de letter. Hiërarchie mag uit gewicht en kleur komen in plaats "
                        "van uit de maat, maar hier doet geen van beide het werk: draagt "
                        "de vulling of de compositie het, kijk dan op de render of dat "
                        "werkt, en zo niet, dan loopt deze slide op één maat.")

    # --- deckbrede tellingen ------------------------------------------------
    rollen_uit = {}
    for rol, maten in sorted(maten_per_rol.items()):
        rollen_uit[rol] = {str(pt): sorted(set(slides))
                           for pt, slides in sorted(maten.items())}
        if len(maten) > 1 and rol not in BAND_ROLLEN:
            reeks = ", ".join(f"{pt:.0f}pt op slide {sorted(set(s))}"
                              for pt, s in sorted(maten.items()))
            add(findings, 0, "maten-per-rol",
                f"de rol '{rol}' staat op {len(maten)} maten: {reeks}. Dezelfde rol "
                "houdt deckbreed dezelfde maat (vormentaal §2) — dat zijn de vier "
                "getallen die in de outline zijn vastgelegd. Los is elke slide dan "
                "correct en naast elkaar leest het als net zoveel keer opnieuw beginnen. "
                "In het gemeten deck stond de koprol op 13, 14, 15, 16 en 18pt.",
                "critical")

    body_maten = sorted(maten_per_rol.get("body", {}))
    if len(body_maten) == 1 and body_maten[0] != BODY_NORM:
        add(findings, 0, "bodymaat",
            f"de body staat deckbreed op {body_maten[0]:.0f}pt en de maatstaf zegt "
            f"{BODY_NORM:.0f}pt Lato Light (vormentaal §2). Eén maat per rol is in orde; "
            "de vraag is of deze maat gekozen is. In beide gemeten decks zakte de body "
            "ongemerkt van 16 naar 14pt omdat de tekst er eerst was en de vorm zich "
            "aanpaste.")

    bandruimte = max(1, round(contentslides / 4)) if contentslides else 1
    if len(bandslides) > bandruimte:
        add(findings, 0, "bandfrequentie",
            f"{len(bandslides)} van de {contentslides} contentslides sluiten af met een "
            f"band over de volle breedte (slides {bandslides}), en er is ruimte voor "
            f"{bandruimte} — ten hoogste één per vier slides (vormentaal §10). Een band "
            "van 12,52 bij 1,25 in met één regel erin is dezelfde vorm, of hij navy, "
            "oranje-tint of warmgrijs is; de vulling wisselen doorbreekt de herhaling "
            "niet. De andere slides sluiten af doordat de laatste rij zelf de conclusie "
            "is, doordat één cel vol gekleurd is, of doordat er niets af te sluiten valt.")

    if contentslides and charts + tables == 0 and cijfers >= CIJFERS_VLOER:
        add(findings, 0, "geen-exhibit-bij-cijfers",
            f"nul native grafieken en nul tabellen bij {cijfers} cijferfeiten in de "
            "contentzone. Draagt een deck cijfers, dan zit er minstens één grafiek, "
            "tabel, schema of verdeling in (vormentaal §12), en een financiële reeks van "
            "drie perioden of meer is een tabel — een financiële lezer kan proza niet "
            "vergelijken. Dit is een warn en geen blokkade omdat een schema of een "
            "verdeling met eigen vormen getekend kan zijn: die telt dit script niet, die "
            "ziet alleen de render. Staat er zo'n exhibit, dan is deze melding onjuist; "
            "staan er vier kolommen tekst, dan is hij de kern van de zaak.")

    tellingen = Counter(f["severity"] for f in findings)
    contentwoorden = list(woorden_per_slide.values())

    return {
        "deck": str(deck),
        "slides": len(presentatie.slides._sldIdLst),
        "contentslides": contentslides,
        "charts": charts,
        "tables": tables,
        "cijferfeiten": cijfers,
        # Met severity: de zes tellingen uit het werkplan.
        "maten_per_rol": rollen_uit,
        "bandslides": bandslides,
        "bandruimte": bandruimte,
        "maatsprong": sprongen,
        # Zonder severity: cijfers om te lezen, niet om op af te keuren. Zie de
        # docstring — een drempel hierop leert versnipperen in plaats van reduceren.
        "gemeten_zonder_oordeel": {
            "woorden_per_contentslide": {
                "per_slide": woorden_per_slide,
                "gemiddeld": round(sum(contentwoorden) / len(contentwoorden))
                if contentwoorden else 0,
                "piek": max(contentwoorden) if contentwoorden else 0,
                "ijkpunt": "gemeten op bestaande decks: de voorbeelden 11 tot 14 in "
                           "assets/maatstaf/ staan op 141, 99, 50 en 59 woorden inclusief "
                           "titel, een spreekdeck op gemiddeld 85, en de hoogste meting is "
                           "177 gemiddeld met een piek van 255. Het dichtste voorbeeld is "
                           "tegelijk de sterkste van de vier, dus dit getal is geen "
                           "drempel en ook geen richting — zie de docstring.",
            },
            "woorden_per_element": {
                "elementen": len(woorden_per_element),
                "gemiddeld": round(sum(woorden_per_element) / len(woorden_per_element))
                if woorden_per_element else 0,
                "piek": max(woorden_per_element) if woorden_per_element else 0,
            },
            "plattegronden": [
                {"plattegrond": sleutel, "aantal": aantal,
                 "slides": plattegrond_slides[sleutel]}
                for sleutel, aantal in plattegronden.most_common()
            ],
            "registers": registers(renders) if renders else {
                "gemeten": False,
                "hint": "niet gemeten: de registerverdeling komt uit de render, niet uit "
                        "de XML. Draai render.py en geef de map mee met --renders.",
            },
        },
        "findings": findings,
        "counts": {"critical": tellingen["critical"], "warn": tellingen["warn"]},
        "verdict": ("blocked" if tellingen["critical"]
                    else "clean" if not findings else "warn"),
    }


def _slidenummers(spec: str) -> set[int]:
    """`"12,14-16"` naar `{12, 14, 15, 16}`."""
    uit: set[int] = set()
    for deel in spec.split(","):
        deel = deel.strip()
        if not deel:
            continue
        if "-" in deel:
            van, tot = deel.split("-", 1)
            uit.update(range(int(van), int(tot) + 1))
        else:
            uit.add(int(deel))
    return uit


def scope_nieuw(resultaat: dict, nieuw: set[int]) -> dict:
    """Verplaats bevindingen op slides die je niet zelf hebt gebouwd naar `overgeerfd`.

    Nodig voor de route "een bestaand deck uitbreiden" uit `SKILL.md`. Een deck dat vóór
    deze regels is gebouwd heeft op vrijwel elke slide een gemengde alinea -- op het eerste
    gemeten deck 71 stuks. Zonder deze schifting blokkeert het toevoegen van twee slides op
    tachtig bevindingen in tekst waar je niet aan hebt gezeten, en dan leert de bouwer het
    script te omzeilen in plaats van te gebruiken. De bevindingen verdwijnen niet: ze staan
    apart, want ze zijn waar en iemand mag besluiten het oude deck door te trekken.

    Deckbrede tellingen blijven staan. Die gaan per definitie over het geheel, en een tweede
    maat per rol introduceer je juist wél zelf zodra je een slide toevoegt die niet bij de
    deck past.
    """
    blijft, overgeerfd = [], []
    for f in resultaat["findings"]:
        slide = f.get("slide")
        (blijft if not isinstance(slide, int) or slide in nieuw
         else overgeerfd).append(f)
    resultaat["findings"] = blijft
    resultaat["overgeerfd"] = overgeerfd
    resultaat["nieuw"] = sorted(nieuw)
    resultaat["counts"] = {
        "critical": sum(1 for f in blijft if f["severity"] == "critical"),
        "warn": sum(1 for f in blijft if f["severity"] == "warn"),
        "overgeerfd": len(overgeerfd),
    }
    resultaat["verdict"] = ("blocked" if resultaat["counts"]["critical"]
                           else "clean" if not blijft else "warn")
    return resultaat


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck")
    parser.add_argument("--renders", type=Path, default=None,
                        help="map met de PNG's van dit deck; alleen daarmee is de "
                             "registerverdeling te meten")
    parser.add_argument("--nieuw", default=None, metavar="12,14-16",
                        help="de slides die je zelf hebt gebouwd. Bevindingen op de "
                             "andere slides gaan naar `overgeerfd` en blokkeren niet -- "
                             "voor de route een bestaand deck uitbreiden")
    parser.add_argument("--strict", action="store_true",
                        help="ook een waarschuwing geeft exit 1")
    args = parser.parse_args()

    resultaat = analyse(Path(args.deck), args.renders)
    if args.nieuw:
        resultaat = scope_nieuw(resultaat, _slidenummers(args.nieuw))
    emit(resultaat)

    if resultaat["counts"]["critical"] or (args.strict and resultaat["counts"]["warn"]):
        sys.exit(1)


if __name__ == "__main__":
    main()
